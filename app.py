import streamlit as st
import json
import os
import base64
import hashlib
import tempfile
import time
import re
import uuid

# ════════════════════════════════════════════════════════
# ── BRANDING — edit these lines ──────────────────────────
BRAND_NAME   = "Your Name"           # ← your name or company
BRAND_LOGO   = ""                    # ← path to logo file e.g. "logo.png", or leave ""
APP_SUBTITLE = "Document Intelligence"
# ════════════════════════════════════════════════════════

# ── Pipeline defaults (hidden from end-users) ────────────
DEFAULT_MAX_CHARS   = 3000
DEFAULT_NEW_AFTER   = 2400
DEFAULT_COMBINE     = 500
DEFAULT_TOP_K       = 3
DEFAULT_PERSIST_DIR = "chroma_db"
DEFAULT_EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
DEFAULT_LLM_MODEL   = "models/gemini-2.5-flash"

# ── API keys — reads from Streamlit secrets first, .env fallback ──
try:
    from dotenv import load_dotenv
    load_dotenv()
except Exception:
    pass
try:
    if "GOOGLE_API_KEY" in st.secrets:
        os.environ["GOOGLE_API_KEY"] = st.secrets["GOOGLE_API_KEY"]
    if "GROQ_API_KEY" in st.secrets:
        os.environ["GROQ_API_KEY"] = st.secrets["GROQ_API_KEY"]
except Exception:
    pass

# ── Supabase client ───────────────────────────────────────
try:
    from supabase import create_client, Client as SupabaseClient
    _SUPA_URL = st.secrets.get("SUPABASE_URL", os.getenv("SUPABASE_URL", ""))
    _SUPA_KEY = st.secrets.get("SUPABASE_KEY", os.getenv("SUPABASE_KEY", ""))
    supabase: SupabaseClient = create_client(_SUPA_URL, _SUPA_KEY) if _SUPA_URL and _SUPA_KEY else None
except Exception:
    supabase = None

# ─── Page Config ─────────────────────────────────────────
st.set_page_config(
    page_title=f"{BRAND_NAME} · {APP_SUBTITLE}",
    page_icon="📄",
    layout="centered",
    initial_sidebar_state="collapsed",
)

# ─── CSS ─────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@500;700;800&family=Inter:wght@300;400;500;600&display=swap');

:root {
    --bg:      #0f0f0f;
    --surface: #1a1a1a;
    --surf2:   #222222;
    --border:  #2a2a2a;
    --orange:  #f97316;
    --orange2: #fb923c;
    --ash:     #a3a3a3;
    --white:   #f5f5f5;
    --dim:     #525252;
    --ok:      #22c55e;
    --err:     #ef4444;
}

html, body,
[data-testid="stAppViewContainer"],
[data-testid="stMain"],
[data-testid="block-container"] {
    background-color: var(--bg) !important;
    color: var(--white) !important;
    font-family: 'Inter', sans-serif !important;
}

/* hide sidebar toggle arrow & footer */
[data-testid="collapsedControl"],
[data-testid="stSidebar"],
footer { display: none !important; }

/* ── brand bar ── */
.brand-bar {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 1.4rem 0 1.2rem 0;
    border-bottom: 1px solid var(--border);
    margin-bottom: 2.4rem;
}
.brand-left { display: flex; align-items: center; gap: 13px; }
.brand-logo {
    width: 40px; height: 40px;
    border-radius: 10px;
    object-fit: cover;
    border: 1px solid var(--border);
}
.brand-initials {
    width: 40px; height: 40px;
    background: linear-gradient(135deg, #f97316, #7c2d12);
    border-radius: 10px;
    display: flex; align-items: center; justify-content: center;
    font-size: 1.05rem; font-weight: 800;
    color: white; font-family: 'Syne', sans-serif;
    flex-shrink: 0;
    letter-spacing: -0.02em;
}
.brand-name {
    font-family: 'Syne', sans-serif;
    font-weight: 800; font-size: 1.05rem;
    color: var(--white); letter-spacing: -0.02em;
    line-height: 1.1;
}
.brand-sub {
    font-size: 0.68rem; color: var(--ash);
    letter-spacing: 0.07em; text-transform: uppercase;
}
.brand-pill {
    font-size: 0.62rem; font-weight: 600;
    letter-spacing: 0.08em; text-transform: uppercase;
    padding: 4px 11px; border-radius: 20px;
    border: 1px solid var(--orange); color: var(--orange);
}

/* ── page hero ── */
.hero-title {
    font-family: 'Syne', sans-serif;
    font-size: 2.3rem; font-weight: 800;
    letter-spacing: -0.04em; line-height: 1.1;
    color: var(--white); margin-bottom: 0.4rem;
}
.hero-title span { color: var(--orange); }
.hero-desc {
    color: var(--ash); font-size: 0.88rem;
    line-height: 1.6; margin-bottom: 2rem;
}

/* ── file uploader ── */
[data-testid="stFileUploader"] > div {
    background: var(--surface) !important;
    border: 2px dashed var(--border) !important;
    border-radius: 14px !important;
    transition: border-color .2s;
}
[data-testid="stFileUploader"] > div:hover {
    border-color: var(--orange) !important;
}
[data-testid="stFileUploader"] label { display: none !important; }

/* ── file pill ── */
.file-pill {
    display: flex; align-items: center; gap: 12px;
    background: var(--surface);
    border: 1px solid var(--border);
    border-left: 3px solid var(--orange);
    border-radius: 10px;
    padding: 0.8rem 1.1rem;
    margin: 0.6rem 0 1.2rem 0;
}
.fp-name { font-weight: 600; font-size: 0.88rem; color: var(--white); }
.fp-size { font-size: 0.72rem; color: var(--ash); }

/* ── run button ── */
.stButton > button {
    background: var(--orange) !important;
    color: #0f0f0f !important;
    border: none !important;
    border-radius: 10px !important;
    font-family: 'Syne', sans-serif !important;
    font-weight: 700 !important;
    font-size: 0.93rem !important;
    padding: 0.65rem 1.6rem !important;
    width: 100% !important;
    transition: opacity .15s, transform .15s !important;
}
.stButton > button:hover { opacity:.88 !important; transform:translateY(-1px) !important; }
.stButton > button:disabled { background: var(--surf2) !important; color: var(--dim) !important; }

/* ── metrics row ── */
.metrics-row { display:flex; gap:12px; margin:1.5rem 0; }
.metric-tile {
    flex:1; background:var(--surface);
    border:1px solid var(--border);
    border-radius:10px; padding:1rem; text-align:center;
}
.metric-num {
    font-family:'Syne',sans-serif; font-size:1.8rem;
    font-weight:800; color:var(--orange); line-height:1;
}
.metric-lbl {
    font-size:0.67rem; color:var(--ash);
    text-transform:uppercase; letter-spacing:.08em; margin-top:4px;
}

/* ── divider ── */
.sec-div {
    display:flex; align-items:center; gap:12px;
    margin:2rem 0 1.5rem 0;
}
.sec-div hr { flex:1; border:none; border-top:1px solid var(--border); margin:0; }
.sec-lbl {
    font-size:.68rem; text-transform:uppercase;
    letter-spacing:.1em; color:var(--dim); white-space:nowrap;
}

/* ── answer box ── */
.answer-box {
    background:var(--surface);
    border:1px solid var(--border);
    border-left:3px solid var(--orange);
    border-radius:12px;
    padding:1.3rem 1.5rem;
    margin-bottom: .5rem;
}
/* ensure markdown elements inside answer-box look right */
.answer-box p  { font-size:.9rem; line-height:1.75; color:var(--white); margin:.4rem 0; }
.answer-box li { font-size:.9rem; line-height:1.7;  color:var(--white); }
.answer-box h1,.answer-box h2,.answer-box h3 {
    font-family:'Syne',sans-serif; color:var(--white);
    margin:.8rem 0 .3rem 0;
}
.answer-box code {
    background:var(--surf2); border:1px solid var(--border);
    border-radius:4px; padding:1px 5px;
    font-size:.82rem; color:var(--orange2);
}
.answer-box pre {
    background:var(--surf2); border:1px solid var(--border);
    border-radius:8px; padding:1rem;
    overflow-x:auto; font-size:.8rem;
}
.answer-box table {
    width:100%; border-collapse:collapse;
    font-size:.83rem; margin:.6rem 0;
}
.answer-box th {
    background:var(--surf2); color:var(--orange2);
    padding:.5rem .75rem; text-align:left;
    border:1px solid var(--border); font-weight:600;
}
.answer-box td {
    padding:.45rem .75rem; border:1px solid var(--border);
    color:var(--ash);
}
.answer-box tr:nth-child(even) td { background:rgba(255,255,255,.02); }
/* MathJax output */
.answer-box mjx-container { color:var(--white) !important; }

/* ── image label ── */
.img-lbl {
    font-size:.68rem; text-transform:uppercase;
    letter-spacing:.1em; color:var(--ash);
    margin:1rem 0 .5rem 0;
}

/* ── chunk card ── */
.chunk-card {
    background:var(--surf2);
    border:1px solid var(--border);
    border-radius:8px; padding:.85rem 1rem;
    margin-bottom:.55rem; font-size:.77rem;
    line-height:1.5; color:var(--ash);
}

/* ── tags ── */
.tag {
    display:inline-block; padding:1px 7px;
    border-radius:4px; font-size:.59rem;
    font-weight:600; letter-spacing:.06em;
    margin-right:4px; text-transform:uppercase;
}
.t-text  { background:rgba(249,115,22,.1); color:#fb923c; border:1px solid rgba(249,115,22,.2);}
.t-table { background:rgba(163,163,163,.08); color:#a3a3a3; border:1px solid rgba(163,163,163,.18);}
.t-image { background:rgba(34,197,94,.08); color:#4ade80; border:1px solid rgba(34,197,94,.18);}

/* ── chat ── */
[data-testid="stChatMessage"] { background:transparent !important; }
[data-testid="stChatMessageContent"] {
    background:var(--surface) !important;
    border:1px solid var(--border) !important;
    border-radius:10px !important;
    color:var(--white) !important;
    font-size:.88rem !important;
}
[data-testid="stChatInput"] textarea {
    background:var(--surface) !important;
    border:1px solid var(--border) !important;
    border-radius:10px !important;
    color:var(--white) !important;
    font-size:.88rem !important;
}
[data-testid="stChatInput"] textarea:focus {
    border-color:var(--orange) !important;
    box-shadow:0 0 0 2px rgba(249,115,22,.12) !important;
}

/* ── expander ── */
[data-testid="stExpander"] {
    background:var(--surf2) !important;
    border:1px solid var(--border) !important;
    border-radius:10px !important;
}
[data-testid="stExpander"] summary {
    font-size:.79rem !important; color:var(--ash) !important;
}

/* ── progress ── */
.stProgress > div > div > div { background:var(--orange) !important; }

/* ── tabs ── */
[data-testid="stTabs"] [role="tablist"] {
    gap:4px; border-bottom:1px solid var(--border) !important;
}
[data-testid="stTabs"] button[role="tab"] {
    font-family:'Inter',sans-serif !important;
    font-size:.8rem !important; font-weight:500 !important;
    color:var(--ash) !important; background:transparent !important;
    border-radius:6px 6px 0 0 !important; padding:6px 16px !important;
}
[data-testid="stTabs"] button[role="tab"][aria-selected="true"] {
    color:var(--orange) !important;
    border-bottom:2px solid var(--orange) !important;
    background:rgba(249,115,22,.06) !important;
}

/* ── log panel ── */
.log-wrap {
    background:var(--surface); border:1px solid var(--border);
    border-radius:10px; padding:1rem 1.2rem;
    font-family:'Courier New',monospace; font-size:.74rem;
    max-height:360px; overflow-y:auto;
}
.ll { padding:2px 0; border-bottom:1px solid var(--border); color:var(--dim); }
.ll.ok  { color:var(--ok); }
.ll.err { color:var(--err); }
.ll .px { color:var(--orange2); margin-right:6px; }

hr { border-color:var(--border) !important; }
label { color:var(--ash) !important; font-size:.75rem !important; }

/* ── quiz ── */
.quiz-q {
    background:var(--surface);
    border:1px solid var(--border);
    border-radius:12px;
    padding:1.3rem 1.5rem;
    margin-bottom:1.2rem;
}
.quiz-q-num {
    font-size:.65rem; text-transform:uppercase;
    letter-spacing:.1em; color:var(--orange);
    font-weight:600; margin-bottom:.4rem;
}
.quiz-q-text {
    font-size:.95rem; font-weight:500;
    color:var(--white); line-height:1.5;
    margin-bottom:1rem;
}
.quiz-opt {
    display:flex; align-items:center; gap:10px;
    background:var(--surf2); border:1px solid var(--border);
    border-radius:8px; padding:.65rem 1rem;
    margin-bottom:.5rem; cursor:pointer;
    font-size:.85rem; color:var(--ash);
    transition:border-color .15s, color .15s;
}
.quiz-opt:hover { border-color:var(--orange); color:var(--white); }
.quiz-opt.correct {
    border-color:#22c55e; color:#22c55e;
    background:rgba(34,197,94,.08);
}
.quiz-opt.wrong {
    border-color:#ef4444; color:#ef4444;
    background:rgba(239,68,68,.08);
}
.quiz-score-card {
    background:var(--surface);
    border:1px solid var(--border);
    border-radius:14px; padding:2rem;
    text-align:center; margin-bottom:1.5rem;
}
.quiz-score-num {
    font-family:'Syne',sans-serif;
    font-size:3.5rem; font-weight:800;
    color:var(--orange); line-height:1;
}
.quiz-score-lbl {
    font-size:.75rem; color:var(--ash);
    text-transform:uppercase; letter-spacing:.1em;
    margin-top:.3rem;
}
.quiz-explanation {
    background:var(--surf2);
    border-left:3px solid var(--orange);
    border-radius:0 8px 8px 0;
    padding:.7rem 1rem;
    font-size:.8rem; color:var(--ash);
    margin-top:.5rem; line-height:1.5;
}
.diff-badge {
    display:inline-block; padding:3px 10px;
    border-radius:20px; font-size:.65rem;
    font-weight:700; letter-spacing:.07em;
    text-transform:uppercase; margin-right:6px;
}
.diff-easy   { background:rgba(34,197,94,.12); color:#4ade80; border:1px solid rgba(34,197,94,.25);}
.diff-medium { background:rgba(249,115,22,.12); color:#fb923c; border:1px solid rgba(249,115,22,.25);}
.diff-hard   { background:rgba(239,68,68,.12);  color:#f87171; border:1px solid rgba(239,68,68,.25);}

/* ── general knowledge answer box ── */
.gk-banner {
    background: rgba(245,158,11,.08);
    border: 1px solid rgba(245,158,11,.25);
    border-radius: 10px;
    padding: .6rem 1rem;
    display: flex;
    align-items: flex-start;
    gap: 10px;
    margin-bottom: .75rem;
}
.gk-banner-icon { font-size: 1.1rem; flex-shrink:0; margin-top:1px; }
.gk-banner-text { font-size: .78rem; color: #fcd34d; line-height: 1.5; }
.gk-banner-text strong { color: #fbbf24; }
.answer-box-gk {
    background: var(--surface);
    border: 1px solid rgba(245,158,11,.3);
    border-left: 3px solid #f59e0b;
    border-radius: 12px;
    padding: 1.3rem 1.5rem;
    margin-bottom: .5rem;
}
.answer-box-gk p  { font-size:.9rem; line-height:1.75; color:var(--white); margin:.4rem 0; }
.answer-box-gk li { font-size:.9rem; line-height:1.7;  color:var(--white); }
.answer-box-gk h1,.answer-box-gk h2,.answer-box-gk h3 {
    font-family:'Syne',sans-serif; color:var(--white); margin:.8rem 0 .3rem 0;
}
.answer-box-gk code {
    background:var(--surf2); border:1px solid var(--border);
    border-radius:4px; padding:1px 5px;
    font-size:.82rem; color:#fbbf24;
}
.answer-box-gk table { width:100%; border-collapse:collapse; font-size:.83rem; margin:.6rem 0; }
.answer-box-gk th   { background:var(--surf2); color:#fbbf24; padding:.5rem .75rem; text-align:left; border:1px solid var(--border); }
.answer-box-gk td   { padding:.45rem .75rem; border:1px solid var(--border); color:var(--ash); }

/* ── summary ── */
.sum-hero {
    background: linear-gradient(135deg, #1a1a1a 0%, #222 100%);
    border: 1px solid var(--border);
    border-left: 4px solid var(--orange);
    border-radius: 14px;
    padding: 1.8rem 2rem;
    margin-bottom: 1.5rem;
}
.sum-doc-title {
    font-family:'Syne',sans-serif;
    font-weight:800; font-size:1.5rem;
    color:var(--white); letter-spacing:-.02em;
    margin-bottom:.3rem;
}
.sum-doc-type {
    font-size:.72rem; color:var(--orange);
    text-transform:uppercase; letter-spacing:.1em;
    font-weight:600; margin-bottom:.8rem;
}
.sum-plain-english {
    font-size:.95rem; color:#d4d4d4;
    line-height:1.8; border-top:1px solid var(--border);
    padding-top:.9rem; margin-top:.6rem;
}
.sum-section {
    background:var(--surface);
    border:1px solid var(--border);
    border-radius:12px;
    padding:1.4rem 1.6rem;
    margin-bottom:1.1rem;
}
.sum-section-title {
    font-family:'Syne',sans-serif;
    font-weight:700; font-size:1rem;
    color:var(--orange); margin-bottom:.9rem;
    display:flex; align-items:center; gap:8px;
}
.sum-section p  { font-size:.88rem; line-height:1.8; color:#d4d4d4; margin:.3rem 0; }
.sum-section li { font-size:.88rem; line-height:1.75; color:#d4d4d4; }
.sum-section h3 { font-family:'Syne',sans-serif; color:var(--white);
                  font-size:.95rem; margin:.8rem 0 .3rem 0; }
.sum-section code {
    background:var(--surf2); border:1px solid var(--border);
    border-radius:4px; padding:1px 6px;
    font-size:.82rem; color:var(--orange2);
}
.sum-section table { width:100%; border-collapse:collapse; font-size:.83rem; margin:.5rem 0; }
.sum-section th { background:var(--surf2); color:var(--orange2);
                  padding:.5rem .75rem; text-align:left; border:1px solid var(--border); }
.sum-section td { padding:.45rem .75rem; border:1px solid var(--border); color:var(--ash); }
.sum-section tr:nth-child(even) td { background:rgba(255,255,255,.02); }
.sum-concept-card {
    background:var(--surf2);
    border:1px solid var(--border);
    border-left:3px solid var(--orange);
    border-radius:8px;
    padding:.8rem 1rem;
    margin-bottom:.6rem;
}
.sum-concept-term {
    font-weight:700; font-size:.88rem;
    color:var(--orange2); margin-bottom:.25rem;
}
.sum-concept-def { font-size:.83rem; color:var(--ash); line-height:1.6; }
.sum-takeaway {
    background: rgba(249,115,22,.07);
    border:1px solid rgba(249,115,22,.2);
    border-radius:8px; padding:.7rem 1rem;
    margin-bottom:.5rem; font-size:.87rem;
    color:var(--white); line-height:1.6;
    display:flex; gap:10px; align-items:flex-start;
}
.sum-takeaway-num {
    background:var(--orange); color:#000;
    font-weight:800; font-size:.7rem;
    border-radius:50%; width:20px; height:20px;
    display:flex; align-items:center; justify-content:center;
    flex-shrink:0; margin-top:1px;
}
.sum-fig-caption {
    font-size:.72rem; color:var(--ash);
    text-align:center; margin-top:.4rem;
    font-style:italic;
}
.sum-regen-bar {
    display:flex; justify-content:flex-end;
    margin-bottom:1rem;
}
/* ── auth screens ── */
.auth-wrap {
    max-width: 420px;
    margin: 3rem auto 0 auto;
}
.auth-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 16px;
    padding: 2.4rem 2.2rem 2rem 2.2rem;
}
.auth-logo {
    text-align: center;
    margin-bottom: 1.6rem;
}
.auth-logo-circle {
    width: 52px; height: 52px;
    background: var(--orange);
    border-radius: 50%;
    display: inline-flex;
    align-items: center;
    justify-content: center;
    font-family: 'Syne', sans-serif;
    font-weight: 800; font-size: 1.3rem;
    color: #000; margin-bottom: .5rem;
}
.auth-title {
    font-family: 'Syne', sans-serif;
    font-weight: 800; font-size: 1.35rem;
    color: var(--white); text-align: center;
    margin-bottom: .3rem;
}
.auth-sub {
    font-size: .8rem; color: var(--ash);
    text-align: center; margin-bottom: 1.6rem;
}
.auth-divider {
    display: flex; align-items: center;
    gap: .75rem; margin: 1rem 0;
}
.auth-divider hr {
    flex: 1; border: none;
    border-top: 1px solid var(--border);
}
.auth-divider span {
    font-size: .72rem; color: #525252;
    white-space: nowrap;
}
.auth-switch {
    text-align: center;
    font-size: .8rem; color: var(--ash);
    margin-top: 1.2rem;
}
.auth-switch strong { color: var(--orange2); cursor: pointer; }
.auth-err {
    background: rgba(239,68,68,.08);
    border: 1px solid rgba(239,68,68,.25);
    border-radius: 8px; padding: .6rem .9rem;
    font-size: .8rem; color: #fca5a5;
    margin-bottom: .8rem;
}
.auth-ok {
    background: rgba(34,197,94,.08);
    border: 1px solid rgba(34,197,94,.25);
    border-radius: 8px; padding: .6rem .9rem;
    font-size: .8rem; color: #86efac;
    margin-bottom: .8rem;
}
/* user pill in header */
.user-pill {
    display: inline-flex; align-items: center;
    gap: 7px; background: var(--surf2);
    border: 1px solid var(--border);
    border-radius: 999px;
    padding: .3rem .75rem .3rem .4rem;
    font-size: .75rem; color: var(--ash);
}
.user-pill-dot {
    width: 8px; height: 8px;
    background: #22c55e; border-radius: 50%;
    flex-shrink: 0;
}
/* doc selector card */
.doc-card {
    background: var(--surf2);
    border: 1px solid var(--border);
    border-radius: 10px;
    padding: .75rem 1rem;
    margin-bottom: .5rem;
    cursor: pointer;
    transition: border-color .15s;
}
.doc-card:hover  { border-color: var(--orange); }
.doc-card.active { border-color: var(--orange); background: rgba(249,115,22,.06); }
.doc-card-name   { font-size: .88rem; font-weight: 600; color: var(--white); }
.doc-card-meta   { font-size: .72rem; color: var(--ash); margin-top: 2px; }
/* ── premium / waitlist ── */
.premium-badge {
    display: inline-flex; align-items: center; gap: 5px;
    background: linear-gradient(135deg, #f97316, #fb923c);
    color: #000; font-size: .65rem; font-weight: 800;
    font-family: 'Syne', sans-serif;
    letter-spacing: .06em; text-transform: uppercase;
    padding: .2rem .55rem; border-radius: 999px;
    cursor: pointer;
}
.upgrade-nudge {
    background: linear-gradient(135deg, rgba(249,115,22,.1), rgba(251,146,60,.05));
    border: 1px solid rgba(249,115,22,.3);
    border-radius: 12px; padding: 1.1rem 1.4rem;
    margin: 1rem 0; display: flex;
    align-items: center; justify-content: space-between;
    gap: 1rem; flex-wrap: wrap;
}
.upgrade-nudge-text { flex: 1; min-width: 180px; }
.upgrade-nudge-title {
    font-family: 'Syne', sans-serif; font-weight: 700;
    font-size: .92rem; color: #fff; margin-bottom: .2rem;
}
.upgrade-nudge-sub { font-size: .78rem; color: #a3a3a3; line-height: 1.5; }
.waitlist-hero {
    text-align: center; padding: 2rem 1rem 1.5rem 1rem;
}
.waitlist-hero-icon { font-size: 2.5rem; margin-bottom: .6rem; }
.waitlist-hero-title {
    font-family: 'Syne', sans-serif; font-weight: 800;
    font-size: 1.5rem; color: #fff;
    margin-bottom: .4rem;
}
.waitlist-hero-sub { font-size: .88rem; color: #a3a3a3; line-height: 1.7; }
.feature-row {
    display: flex; gap: .6rem; align-items: flex-start;
    padding: .65rem 0;
    border-bottom: 1px solid #1f1f1f;
}
.feature-row:last-child { border-bottom: none; }
.feature-icon { font-size: 1.1rem; flex-shrink: 0; margin-top: 1px; }
.feature-text { font-size: .85rem; color: #d4d4d4; line-height: 1.5; }
.feature-text strong { color: #fff; }
.waitlist-joined {
    background: rgba(34,197,94,.08);
    border: 1px solid rgba(34,197,94,.25);
    border-radius: 12px; padding: 1.4rem;
    text-align: center; margin: 1rem 0;
}
.waitlist-joined-title {
    font-family: 'Syne', sans-serif; font-weight: 800;
    font-size: 1.1rem; color: #4ade80; margin-bottom: .4rem;
}
.waitlist-joined-sub { font-size: .82rem; color: #a3a3a3; }
/* ── queue / busy card ── */
.busy-card {
    background:var(--surface);
    border:1px solid var(--border);
    border-left:3px solid var(--orange);
    border-radius:14px;
    padding:2rem 1.8rem;
    text-align:center;
    margin-top:1rem;
}
.busy-spinner {
    width:40px; height:40px;
    border:3px solid var(--border);
    border-top:3px solid var(--orange);
    border-radius:50%;
    animation:spin 1s linear infinite;
    margin:0 auto 1.2rem auto;
}
@keyframes spin { to { transform:rotate(360deg); } }
.busy-title {
    font-family:'Syne',sans-serif;
    font-weight:800; font-size:1.15rem;
    color:var(--white); margin-bottom:.4rem;
}
.busy-sub { color:var(--ash); font-size:.83rem; line-height:1.6; }
.busy-sub strong { color:var(--orange2); }
</style>
""", unsafe_allow_html=True)

# ── Auth + session state ──────────────────────────────────
for k, v in {
    # ── auth ──
    "user": None,              # Supabase user object after sign-in
    "profile": None,           # row from public.profiles
    "auth_screen": "signin",   # "signin" | "signup" | "reset"
    "auth_error": "",
    "auth_ok": "",
    "active_doc_id": None,     # UUID of currently loaded document
    # ── pipeline ──
    "db": None,
    "processed_chunks": [],
    "pipeline_ran": False,
    "logs": [],
    "metrics": {"elements": 0, "chunks": 0, "docs": 0},
    "chat_history": [],
    "pipeline_busy": False,
    "all_page_images": {},
    # ── summary ──
    "summary": None,
    "summary_images": [],
    "summary_tables": [],
    "doc_name": "",
    # ── waitlist ──
    "on_waitlist": False,       # True once user joins
    "show_upgrade": False,      # True to jump to upgrade tab
    "upgrade_trigger": "",      # which trigger surfaced the nudge
    "quiz_questions": [],
    "quiz_answers": {},
    "quiz_submitted": False,
    "quiz_generating": False,
}.items():
    if k not in st.session_state:
        st.session_state[k] = v

# Derive user_id: use Supabase uid when signed in, else temp UUID
if st.session_state.user:
    _uid = st.session_state.user.id
else:
    if "anon_uid" not in st.session_state:
        st.session_state.anon_uid = uuid.uuid4().hex
    _uid = st.session_state.anon_uid

USER_PERSIST_DIR = os.path.join(DEFAULT_PERSIST_DIR, _uid)

# ── Global busy flag (shared across all sessions via a temp file) ──────────
# Stores JSON: {"user_id": "...", "started_at": unix_timestamp}
# Auto-expires after PIPELINE_TIMEOUT_SECS so a crashed session never
# blocks everyone permanently.
_BUSY_FLAG_PATH      = os.path.join(tempfile.gettempdir(), "rag_pipeline_busy.flag")
PIPELINE_TIMEOUT_SECS = 600   # 10 minutes — more than enough for any doc

def _flag_data() -> dict:
    """Read the flag file. Returns {} if missing or corrupt."""
    try:
        with open(_BUSY_FLAG_PATH) as f:
            return json.loads(f.read())
    except Exception:
        return {}

def _global_busy() -> bool:
    """True if a pipeline is running AND hasn't timed out."""
    d = _flag_data()
    if not d:
        return False
    elapsed = time.time() - d.get("started_at", 0)
    if elapsed > PIPELINE_TIMEOUT_SECS:
        _set_global_busy(False)   # auto-expire
        return False
    return True

def _set_global_busy(busy: bool, user_id: str = ""):
    if busy:
        with open(_BUSY_FLAG_PATH, "w") as f:
            json.dump({"user_id": user_id, "started_at": time.time()}, f)
    else:
        try:
            os.remove(_BUSY_FLAG_PATH)
        except FileNotFoundError:
            pass

def _busy_user() -> str:
    return _flag_data().get("user_id", "")

def _busy_elapsed_secs() -> int:
    d = _flag_data()
    if not d:
        return 0
    return int(time.time() - d.get("started_at", time.time()))


# ─── Supabase DB helpers ─────────────────────────────────

def db_get_profile(user_id: str) -> dict:
    """Fetch profile row for a user."""
    try:
        r = supabase.table("profiles").select("*").eq("id", user_id).single().execute()
        return r.data or {}
    except Exception:
        return {}

def db_upsert_profile(user_id: str, full_name: str, email: str):
    """Create or update profile."""
    try:
        supabase.table("profiles").upsert({
            "id": user_id, "full_name": full_name,
            "email": email, "last_seen": "now()"
        }).execute()
    except Exception:
        pass

def db_touch_last_seen(user_id: str):
    try:
        supabase.table("profiles").update({"last_seen": "now()"}).eq("id", user_id).execute()
    except Exception:
        pass

def db_save_document(user_id: str, name: str, file_type: str,
                     chunk_count: int, page_count: int, persist_dir: str) -> str:
    """Insert a document record, return its UUID."""
    try:
        r = supabase.table("documents").insert({
            "user_id":     user_id,
            "name":        name,
            "file_type":   file_type,
            "chunk_count": chunk_count,
            "page_count":  page_count,
            "persist_dir": persist_dir,
        }).execute()
        return r.data[0]["id"] if r.data else None
    except Exception:
        return None

def db_get_documents(user_id: str) -> list:
    """Return all documents for a user, newest first."""
    try:
        r = supabase.table("documents").select("*") \
            .eq("user_id", user_id).order("created_at", desc=True).execute()
        return r.data or []
    except Exception:
        return []

def db_delete_document(doc_id: str):
    try:
        supabase.table("documents").delete().eq("id", doc_id).execute()
    except Exception:
        pass

def db_save_chat_session(user_id: str, document_id: str) -> str:
    """Create a chat session, return its UUID."""
    try:
        r = supabase.table("chat_sessions").insert({
            "user_id": user_id, "document_id": document_id
        }).execute()
        return r.data[0]["id"] if r.data else None
    except Exception:
        return None

def db_save_message(session_id: str, role: str, content: str, answer_type: str = "doc"):
    try:
        supabase.table("chat_messages").insert({
            "session_id":  session_id,
            "role":        role,
            "content":     content,
            "answer_type": answer_type,
        }).execute()
    except Exception:
        pass

def db_get_messages(session_id: str) -> list:
    try:
        r = supabase.table("chat_messages").select("*") \
            .eq("session_id", session_id).order("created_at").execute()
        return r.data or []
    except Exception:
        return []

def db_save_quiz(user_id: str, document_id: str,
                 difficulty: str, score: int, total: int):
    try:
        supabase.table("quiz_sessions").insert({
            "user_id":     user_id,
            "document_id": document_id,
            "difficulty":  difficulty,
            "score":       score,
            "total":       total,
        }).execute()
    except Exception:
        pass

def db_join_waitlist(user_id: str, full_name: str, email: str,
                     trigger: str, price_range: str) -> bool:
    """Add user to waitlist. Returns True on success."""
    try:
        supabase.table("waitlist").insert({
            "user_id":     user_id,
            "full_name":   full_name,
            "email":       email,
            "trigger":     trigger,
            "price_range": price_range,
        }).execute()
        return True
    except Exception:
        return False

def db_is_on_waitlist(user_id: str) -> bool:
    """Check if user already joined the waitlist."""
    try:
        r = supabase.table("waitlist").select("id") \
            .eq("user_id", user_id).limit(1).execute()
        return bool(r.data)
    except Exception:
        return False

def db_waitlist_count() -> int:
    """Total number of waitlist signups — shown as social proof."""
    try:
        r = supabase.table("waitlist").select("id", count="exact").execute()
        return r.count or 0
    except Exception:
        return 0

# ─── Helpers ─────────────────────────────────────────────
def log(msg, level="info"):
    st.session_state.logs.append({"msg": msg, "level": level})

def _hash(s: str) -> str:
    """MD5 fingerprint of a string — used for deduplication."""
    return hashlib.md5(s.encode()).hexdigest()

def collect_content(chunks):
    """
    Walk retrieved chunks and return three deduplicated lists:
      - unique_images  : base64 strings (no duplicates by content hash)
      - unique_tables  : HTML strings   (no duplicates by content hash)
      - unique_texts   : raw text blocks (no duplicates by content hash)
    """
    seen_imgs   = set()
    seen_tables = set()
    seen_texts  = set()

    unique_images  = []
    unique_tables  = []
    unique_texts   = []

    for c in chunks:
        if "original_content" not in c.metadata:
            continue
        orig = json.loads(c.metadata["original_content"])

        # ── images ──
        for b64 in orig.get("images_base64", []):
            if not b64:
                continue
            h = _hash(b64)
            if h not in seen_imgs:
                seen_imgs.add(h)
                unique_images.append(b64)

        # ── tables ──
        for tbl in orig.get("tables_html", []):
            if not tbl:
                continue
            h = _hash(tbl.strip())
            if h not in seen_tables:
                seen_tables.add(h)
                unique_tables.append(tbl)

        # ── raw text ──
        txt = orig.get("raw_text", "").strip()
        if txt:
            h = _hash(txt)
            if h not in seen_texts:
                seen_texts.add(h)
                unique_texts.append(txt)

    return unique_images, unique_tables, unique_texts


# ── Image extraction ──────────────────────────────────────────────────────────

def extract_images_from_pdf(pdf_path: str, dpi: int = 150) -> dict:
    """
    Render every page of a PDF as a PNG at `dpi` resolution using PyMuPDF.
    Returns {page_number (1-based): base64_png_string}.
    Captures everything — raster images, vector graphics, diagrams, formulas.
    """
    try:
        import fitz  # PyMuPDF
        doc    = fitz.open(pdf_path)
        pages  = {}
        matrix = fitz.Matrix(dpi / 72, dpi / 72)   # 72 is PDF's default DPI
        for page in doc:
            pix  = page.get_pixmap(matrix=matrix, alpha=False)
            b64  = base64.b64encode(pix.tobytes("png")).decode()
            pages[page.number + 1] = b64           # 1-based page numbers
        doc.close()
        return pages
    except Exception as e:
        return {}


def extract_images_from_docx(docx_path: str) -> list:
    """Extract embedded images from a Word document as base64 PNGs."""
    try:
        from docx import Document as DocxDocument
        doc    = DocxDocument(docx_path)
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                img_bytes = rel.target_part.blob
                images.append(base64.b64encode(img_bytes).decode())
        return images
    except Exception:
        return []


def extract_images_from_pptx(pptx_path: str) -> list:
    """Extract one rendered image per slide from a PowerPoint file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs    = Presentation(pptx_path)
        images = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:          # MSO_SHAPE_TYPE.PICTURE = 13
                    img_bytes = shape.image.blob
                    images.append(base64.b64encode(img_bytes).decode())
        return images
    except Exception:
        return []


def attach_page_images_to_chunks(chunks: list, page_images: dict) -> list:
    """
    For each chunk, find which pages its elements came from and attach
    the corresponding page renders into the chunk's image list.
    Stores result in chunk._page_nums (a set of ints).
    """
    for chunk in chunks:
        page_nums = set()
        if hasattr(chunk, "metadata"):
            if hasattr(chunk.metadata, "page_number") and chunk.metadata.page_number:
                page_nums.add(int(chunk.metadata.page_number))
            if hasattr(chunk.metadata, "orig_elements"):
                for el in chunk.metadata.orig_elements:
                    if hasattr(el, "metadata") and hasattr(el.metadata, "page_number"):
                        if el.metadata.page_number:
                            page_nums.add(int(el.metadata.page_number))
        try:
            chunk._page_nums = page_nums
        except Exception:
            pass   # some unstructured objects are frozen — we'll fall back below
    return chunks


# ── General knowledge detection ───────────────────────────────────────────────

# Cosine distance threshold — Chroma returns distances (lower = more similar).
# Above this threshold means the document doesn't have a good answer.
GK_DISTANCE_THRESHOLD = 0.65

# Regex patterns that signal the user wants an explanation, not a lookup
GK_INTENT_PATTERNS = re.compile(
    r"\b(what\s+is|what\s+are|what\s+does|explain|define|definition\s+of|"
    r"how\s+does|how\s+do|tell\s+me\s+about|describe|meaning\s+of|"
    r"what\s+do\s+you\s+mean|elaborate|break\s+it\s+down|"
    r"i\s+don.?t\s+understand|help\s+me\s+understand)\b",
    re.IGNORECASE
)

def should_use_general_knowledge(query: str, retrieved_docs, scores: list) -> tuple:
    """
    Returns (use_gk: bool, reason: str).
    use_gk is True when:
      - best similarity score is too low (doc doesn't contain the answer), OR
      - query intent is clearly explanatory (what is X, explain Y, etc.)
    """
    low_confidence = scores and min(scores) > GK_DISTANCE_THRESHOLD
    explanatory    = bool(GK_INTENT_PATTERNS.search(query))

    if low_confidence and explanatory:
        return True, "both"
    if low_confidence:
        return True, "low_confidence"
    if explanatory:
        return True, "intent"
    return False, ""


def render_answer(answer: str, images: list, is_gk: bool = False):
    """
    Render the LLM answer with full markdown + LaTeX support.
    is_gk=True renders with amber styling + disclaimer banner.
    """
    if is_gk:
        st.markdown("""
        <div class="gk-banner">
            <div class="gk-banner-icon">⚠️</div>
            <div class="gk-banner-text">
                <strong>General knowledge answer</strong> — this information comes from the AI's
                training data, <strong>not from your uploaded document</strong>.
                Always verify against your source material.
            </div>
        </div>
        """, unsafe_allow_html=True)
        st.markdown('<div class="answer-box-gk">', unsafe_allow_html=True)
    else:
        st.markdown('<div class="answer-box">', unsafe_allow_html=True)

    st.markdown(answer)
    st.markdown('</div>', unsafe_allow_html=True)

    if images and not is_gk:           # only show doc images for doc answers
        st.markdown('<div class="img-lbl">📷 Referenced figures from document</div>', unsafe_allow_html=True)
        cols = st.columns(min(len(images), 3))
        for i, b64 in enumerate(images[:6]):
            try:
                cols[i % len(cols)].image(base64.b64decode(b64), use_container_width=True)
            except Exception:
                pass

def _is_quota_error(e: Exception) -> bool:
    s = str(e)
    return "429" in s or "RESOURCE_EXHAUSTED" in s or "quota" in s.lower()

def _groq_messages(gemini_messages):
    """
    Convert a LangChain HumanMessage list (which may contain image_url blocks)
    into plain-text-only messages safe for Groq (no vision support).
    Images are dropped; everything else is kept.
    """
    from langchain_core.messages import HumanMessage as HM
    plain_parts = []
    for msg in gemini_messages:
        if hasattr(msg, "content"):
            if isinstance(msg.content, str):
                plain_parts.append(msg.content)
            elif isinstance(msg.content, list):
                for block in msg.content:
                    if isinstance(block, dict) and block.get("type") == "text":
                        plain_parts.append(block["text"])
                    # image_url blocks are silently dropped — Groq can't handle them
    return [HM(content="\n\n".join(plain_parts))]

def invoke_with_fallback(messages, status_slot=None):
    """
    Tier 1 — Gemini (vision-capable, 20 req/day free).
    Tier 2 — Groq / llama-3.3-70b (text+tables only, 14,400 req/day free).

    On a quota/rate-limit error from Gemini, switches to Groq immediately
    and shows the user a small notice. Any non-quota error is re-raised.
    """
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_core.messages import HumanMessage

    # ── Tier 1: Gemini ───────────────────────────────────
    try:
        llm = ChatGoogleGenerativeAI(model=DEFAULT_LLM_MODEL, temperature=0)
        return llm.invoke(messages), "gemini"
    except Exception as e:
        if not _is_quota_error(e):
            raise
        # quota hit → fall through to Groq
        if status_slot:
            status_slot.markdown(
                '<div style="background:#1a1a1a; border:1px solid #2a2a2a; '
                'border-left:3px solid #f97316; border-radius:10px; '
                'padding:0.75rem 1.1rem; font-size:0.8rem; color:#a3a3a3; margin-bottom:8px;">'
                '⚡ Gemini quota reached — switching to <strong style="color:#f97316">Groq (Llama 3.3)</strong>. '
                'Images won\'t be analysed this turn but will still display.</div>',
                unsafe_allow_html=True
            )

    # ── Tier 2: Groq ────────────────────────────────────
    try:
        from langchain_groq import ChatGroq
        groq_llm  = ChatGroq(model="llama-3.3-70b-versatile", temperature=0)
        groq_msgs = _groq_messages(messages)   # strip image blocks
        return groq_llm.invoke(groq_msgs), "groq"
    except Exception as e2:
        if _is_quota_error(e2):
            raise RuntimeError(
                "Both Gemini and Groq have hit their rate limits. "
                "Please wait a few minutes and try again."
            ) from e2
        raise

def generate_quiz(num_questions: int, difficulty: str) -> list:
    """
    Pull random chunks from the indexed document and ask the LLM to produce
    a JSON quiz. Returns a list of dicts:
      { "question": str, "options": [A,B,C,D], "answer": int (0-3), "explanation": str }
    """
    from langchain_core.messages import HumanMessage

    # sample up to 6 chunks as context so we cover the doc breadth
    docs = st.session_state.processed_chunks
    import random
    sample = random.sample(docs, min(6, len(docs)))
    context = "\n\n---\n\n".join(
        json.loads(d.metadata["original_content"]).get("raw_text", d.page_content)[:600]
        for d in sample
    )

    diff_instruction = {
        "Easy":   "straightforward recall questions — answers are stated directly in the text",
        "Medium": "questions that require understanding and connecting ideas across the text",
        "Hard":   "challenging questions requiring inference, analysis, or application of concepts",
    }[difficulty]

    prompt = f"""You are a quiz generator. Based on the document excerpts below, create exactly {num_questions} multiple-choice questions.

Difficulty: {difficulty} — {diff_instruction}

DOCUMENT EXCERPTS:
{context}

STRICT OUTPUT FORMAT — respond with valid JSON only, no markdown fences, no extra text:
[
  {{
    "question": "Question text here?",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "answer": 0,
    "explanation": "Brief explanation of why the correct answer is right."
  }}
]

Rules:
- "answer" is the 0-based index of the correct option (0=A, 1=B, 2=C, 3=D)
- All 4 options must be plausible — no obviously wrong distractors
- Questions must be answerable from the document content provided
- Return exactly {num_questions} question objects in the array
- No markdown, no code fences — raw JSON only"""

    msgs = [HumanMessage(content=prompt)]
    response, _ = invoke_with_fallback(msgs)
    raw = response.content.strip()

    # strip any accidental markdown fences
    raw = re.sub(r"^```[a-z]*\n?", "", raw, flags=re.MULTILINE)
    raw = re.sub(r"```$", "", raw, flags=re.MULTILINE).strip()

    questions = json.loads(raw)

    # validate shape
    for q in questions:
        assert "question" in q and "options" in q and "answer" in q
        assert len(q["options"]) == 4
        assert 0 <= q["answer"] <= 3

    return questions


def _repair_json(raw: str) -> str:
    """
    Best-effort JSON repair:
    1. Strip anything before the first { and after the last }
    2. Remove markdown fences
    3. If JSON is truncated (common with long LLM output), close any
       open arrays/objects so json.loads can at least parse what's there.
    """
    # remove markdown fences
    raw = re.sub(r"^```[a-zA-Z]*\n?", "", raw, flags=re.MULTILINE)
    raw = re.sub(r"```\s*$",           "", raw, flags=re.MULTILINE).strip()

    # extract from first { to last }
    start = raw.find("{")
    end   = raw.rfind("}")
    if start == -1:
        raise ValueError("No JSON object found in response")
    raw = raw[start: end + 1] if end > start else raw[start:]

    # attempt to close truncated JSON by counting brackets
    try:
        json.loads(raw)
        return raw                      # already valid
    except json.JSONDecodeError:
        pass

    # close any open string literals first (truncated mid-string)
    # count unescaped quotes to see if we're inside a string
    in_string   = False
    escape_next = False
    for ch in raw:
        if escape_next:
            escape_next = False
            continue
        if ch == "\\":
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
    if in_string:
        raw += '"'                      # close the open string

    # close open arrays and objects
    stack = []
    in_str = False
    esc    = False
    for ch in raw:
        if esc:           esc = False;  continue
        if ch == "\\":    esc = True;   continue
        if ch == '"':     in_str = not in_str; continue
        if in_str:        continue
        if ch == "{":     stack.append("}")
        elif ch == "[":   stack.append("]")
        elif ch in "}]" and stack and stack[-1] == ch:
            stack.pop()

    raw += "".join(reversed(stack))    # close in reverse order

    return raw


def _safe_parse_json(raw: str) -> dict:
    """Try to parse JSON, applying repair if needed. Always returns a dict."""
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        try:
            return json.loads(_repair_json(raw))
        except Exception:
            # last resort — return empty structure so UI still renders
            return {}


def _summarise_batch(texts: list, tables: list, images: list,
                     doc_name: str, batch_num: int, total_batches: int) -> dict:
    """
    Summarise one batch of text chunks. Returns a partial summary dict.
    Batching prevents context-window overflow on long documents.
    """
    from langchain_core.messages import HumanMessage

    joined = "\n\n---\n\n".join(texts)
    # hard cap per batch — Groq context limit is ~32k tokens (~24k chars safe)
    if len(joined) > 24_000:
        joined = joined[:24_000] + "\n\n[truncated for length]"

    table_ctx = ""
    for i, t in enumerate(tables[:3]):
        table_ctx += f"\nTABLE {i+1}:\n{t}\n"

    part_note = f"(Part {batch_num} of {total_batches})" if total_batches > 1 else ""

    prompt = f"""You are an expert academic summariser. Summarise the document content below {part_note}.

DOCUMENT: {doc_name}

CONTENT:
{joined}
{f"TABLES:{table_ctx}" if table_ctx else ""}
{f"FIGURES: {len(images)} page image(s) attached." if images else ""}

Return ONLY a valid JSON object — no markdown fences, no explanation, no text outside the JSON.
If a value would contain a double-quote character, escape it as \\".
If a value would contain a newline, use \\n.

JSON STRUCTURE:
{{
  "topic": "One plain sentence — what is this content about",
  "plain_english": "Explain this topic from scratch for a student who has never seen it. Build from basics. Write 3-6 paragraphs. Use \\n\\n between paragraphs.",
  "sections": [
    {{
      "title": "Topic or section name",
      "summary": "Detailed explanation. Use markdown. Use $formula$ for inline math, $$formula$$ for block equations. Cover every idea. Use \\n for line breaks.",
      "formulas": ["$$...$$"],
      "key_point": "The single most important takeaway from this section"
    }}
  ],
  "concepts": [
    {{"term": "Term", "definition": "Simple 1-2 sentence definition"}}
  ],
  "tables": [
    {{"title": "What this table shows", "markdown": "| A | B |\\n|---|---|\\n| x | y |"}}
  ],
  "formulas": [
    {{"label": "Formula name", "latex": "$$...$$", "explanation": "Plain-language meaning"}}
  ],
  "takeaways": ["Key point 1", "Key point 2"]
}}

RULES:
- sections[] must cover EVERY concept in the content — skip nothing real
- Each section summary should be thorough paragraphs, not one-liners
- All math MUST use $ or $$ delimiters
- takeaways[] = 5-10 items minimum
- Return ONLY the JSON object. Nothing else."""

    content = [{"type": "text", "text": prompt}]
    for b64 in images[:2]:
        mime = "image/png" if b64.startswith("iVBOR") else "image/jpeg"
        content.append({"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}})

    response, _ = invoke_with_fallback([HumanMessage(content=content)])
    return _safe_parse_json(response.content)


def _merge_summaries(parts: list) -> dict:
    """Merge multiple partial summary dicts into one coherent summary."""
    if not parts:
        return {}
    if len(parts) == 1:
        return parts[0]

    merged = {
        "topic":         parts[0].get("topic", ""),
        "plain_english": "\n\n".join(p.get("plain_english", "") for p in parts if p.get("plain_english")),
        "sections":      [],
        "concepts":      [],
        "tables":        [],
        "formulas":      [],
        "takeaways":     [],
    }

    seen_sections  = set()
    seen_concepts  = set()
    seen_formulas  = set()
    seen_takeaways = set()

    for p in parts:
        for s in p.get("sections", []):
            key = s.get("title", "").strip().lower()
            if key and key not in seen_sections:
                seen_sections.add(key)
                merged["sections"].append(s)

        for c in p.get("concepts", []):
            key = c.get("term", "").strip().lower()
            if key and key not in seen_concepts:
                seen_concepts.add(key)
                merged["concepts"].append(c)

        for t in p.get("tables", []):
            merged["tables"].append(t)          # tables are always unique

        for f in p.get("formulas", []):
            key = f.get("label", "").strip().lower()
            if key and key not in seen_formulas:
                seen_formulas.add(key)
                merged["formulas"].append(f)

        for t in p.get("takeaways", []):
            key = t.strip().lower()[:60]
            if key and key not in seen_takeaways:
                seen_takeaways.add(key)
                merged["takeaways"].append(t)

    return merged


def generate_summary(doc_name: str) -> dict:
    """
    Generate a full structured summary of the indexed document.

    Strategy:
    - Collect all unique text chunks and tables
    - Split into batches of ~20k chars to stay within Groq's context limit
    - Summarise each batch independently
    - Merge all partial summaries into one coherent result
    - Cache key page images for display
    """
    docs        = st.session_state.processed_chunks
    page_images = st.session_state.get("all_page_images", {})

    # ── Collect unique texts and tables ───────────────────
    all_texts  = []
    all_tables = []
    seen_text_hashes  = set()
    seen_table_hashes = set()

    for doc in docs:
        orig = json.loads(doc.metadata.get("original_content", "{}"))
        txt  = orig.get("raw_text", "").strip()
        if txt:
            h = _hash(txt)
            if h not in seen_text_hashes:
                seen_text_hashes.add(h)
                all_texts.append(txt)
        for tbl in orig.get("tables_html", []):
            if tbl:
                h = _hash(tbl.strip())
                if h not in seen_table_hashes:
                    seen_table_hashes.add(h)
                    all_tables.append(tbl)

    # ── Select key page images (first, last, evenly spaced) ───────────────
    key_images = []
    if page_images:
        sorted_pages = sorted(page_images.keys())
        total        = len(sorted_pages)
        candidates   = [sorted_pages[0]]
        if total > 1:
            step = max(1, total // 5)
            candidates += sorted_pages[1:-1:step]
            candidates.append(sorted_pages[-1])
        candidates = list(dict.fromkeys(candidates))[:6]
        key_images = [page_images[p] for p in candidates]

    # ── Batch texts into ~20k char chunks ─────────────────
    BATCH_CHAR_LIMIT = 20_000
    batches    = []
    current    = []
    current_sz = 0

    for txt in all_texts:
        if current_sz + len(txt) > BATCH_CHAR_LIMIT and current:
            batches.append(current)
            current    = []
            current_sz = 0
        current.append(txt)
        current_sz += len(txt)
    if current:
        batches.append(current)

    if not batches:
        batches = [[""]]

    # ── Summarise each batch ───────────────────────────────
    parts = []
    total_batches = len(batches)
    for i, batch_texts in enumerate(batches):
        # distribute images across batches (2 per batch max)
        batch_imgs = key_images[i * 2: i * 2 + 2]
        # tables only on first batch
        batch_tables = all_tables[:3] if i == 0 else []
        part = _summarise_batch(
            batch_texts, batch_tables, batch_imgs,
            doc_name, i + 1, total_batches
        )
        parts.append(part)

    # ── Merge all parts ────────────────────────────────────
    summary = _merge_summaries(parts)

    # ── Cache images and tables for the UI ────────────────
    st.session_state.summary_images = key_images
    st.session_state.summary_tables = all_tables

    # ── Ensure required keys always exist ─────────────────
    for key, default in {
        "topic": doc_name, "plain_english": "",
        "sections": [], "concepts": [], "tables": [],
        "formulas": [], "takeaways": []
    }.items():
        if key not in summary:
            summary[key] = default

    return summary


# ─── Auth UI ─────────────────────────────────────────────

def _auth_header(title: str, subtitle: str):
    """Renders the logo, brand name, title and subtitle for auth screens."""
    initial = BRAND_NAME[0].upper() if BRAND_NAME else "?"
    st.markdown(f"""
    <div style="text-align:center; padding: 2rem 0 1.5rem 0;">
        <div style="
            width:52px; height:52px; background:#f97316; border-radius:50%;
            display:inline-flex; align-items:center; justify-content:center;
            font-family:'Syne',sans-serif; font-weight:800; font-size:1.3rem;
            color:#000; margin-bottom:.6rem;">
            {initial}
        </div>
        <div style="font-family:'Syne',sans-serif; font-weight:800;
                    font-size:1rem; color:#f5f5f5; margin-bottom:1.2rem;">
            {BRAND_NAME}
        </div>
        <div style="font-family:'Syne',sans-serif; font-weight:800;
                    font-size:1.35rem; color:#ffffff; margin-bottom:.3rem;">
            {title}
        </div>
        <div style="font-size:.82rem; color:#a3a3a3; margin-bottom:.5rem;">
            {subtitle}
        </div>
    </div>
    """, unsafe_allow_html=True)

def render_signin():
    # center the form with columns
    _, col, _ = st.columns([1, 2, 1])
    with col:
        _auth_header("Welcome back", "Sign in to your account")

        if st.session_state.auth_error:
            st.error(st.session_state.auth_error)
        if st.session_state.auth_ok:
            st.success(st.session_state.auth_ok)

        email    = st.text_input("Email address", key="si_email",
                                 placeholder="you@example.com")
        password = st.text_input("Password", type="password", key="si_pass",
                                 placeholder="Your password")

        if st.button("Sign in →", use_container_width=True, type="primary", key="si_btn"):
            st.session_state.auth_error = ""
            st.session_state.auth_ok    = ""
            if not email or not password:
                st.session_state.auth_error = "Please fill in all fields."
                st.rerun()
            else:
                try:
                    res  = supabase.auth.sign_in_with_password(
                        {"email": email, "password": password}
                    )
                    st.session_state.user    = res.user
                    st.session_state.profile = db_get_profile(res.user.id)
                    db_touch_last_seen(res.user.id)
                    st.session_state.auth_error = ""
                    st.session_state.auth_ok    = ""
                    st.rerun()
                except Exception as e:
                    err = str(e)
                    if "Invalid login" in err or "invalid" in err.lower():
                        st.session_state.auth_error = "Incorrect email or password."
                    elif "confirmed" in err.lower() or "verify" in err.lower():
                        st.session_state.auth_error = "Please verify your email — check your inbox."
                    else:
                        st.session_state.auth_error = f"Sign in failed: {err}"
                    st.rerun()

        st.markdown("<div style='height:.5rem'></div>", unsafe_allow_html=True)
        c1, c2 = st.columns(2)
        with c1:
            if st.button("Create account", use_container_width=True, key="si_to_su"):
                st.session_state.auth_screen = "signup"
                st.session_state.auth_error  = ""
                st.session_state.auth_ok     = ""
                st.rerun()
        with c2:
            if st.button("Forgot password", use_container_width=True, key="si_to_rp"):
                st.session_state.auth_screen = "reset"
                st.session_state.auth_error  = ""
                st.session_state.auth_ok     = ""
                st.rerun()


def render_signup():
    _, col, _ = st.columns([1, 2, 1])
    with col:
        _auth_header("Create your account", "Free forever · No credit card needed")

        if st.session_state.auth_error:
            st.error(st.session_state.auth_error)
        if st.session_state.auth_ok:
            st.success(st.session_state.auth_ok)

        full_name = st.text_input("Full name",        key="su_name",
                                  placeholder="Your name")
        email     = st.text_input("Email address",    key="su_email",
                                  placeholder="you@example.com")
        password  = st.text_input("Password",         key="su_pass",
                                  type="password", placeholder="At least 6 characters")
        confirm   = st.text_input("Confirm password", key="su_conf",
                                  type="password", placeholder="Repeat your password")

        if st.button("Create account →", use_container_width=True,
                     type="primary", key="su_btn"):
            st.session_state.auth_error = ""
            st.session_state.auth_ok    = ""
            if not all([full_name, email, password, confirm]):
                st.session_state.auth_error = "Please fill in all fields."
            elif len(password) < 6:
                st.session_state.auth_error = "Password must be at least 6 characters."
            elif password != confirm:
                st.session_state.auth_error = "Passwords do not match."
            else:
                try:
                    res = supabase.auth.sign_up({
                        "email":    email,
                        "password": password,
                        "options":  {"data": {"full_name": full_name}}
                    })
                    if res.user:
                        st.session_state.auth_ok     = (
                            "Account created! Check your email to verify, then sign in."
                        )
                        st.session_state.auth_screen = "signin"
                    else:
                        st.session_state.auth_error = "Sign up failed — please try again."
                except Exception as e:
                    err = str(e)
                    if "already" in err.lower():
                        st.session_state.auth_error = "An account with that email already exists."
                    else:
                        st.session_state.auth_error = f"Sign up failed: {err}"
            st.rerun()

        st.markdown("<div style='height:.5rem'></div>", unsafe_allow_html=True)
        if st.button("← Back to sign in", use_container_width=True, key="su_back"):
            st.session_state.auth_screen = "signin"
            st.session_state.auth_error  = ""
            st.session_state.auth_ok     = ""
            st.rerun()


def render_reset():
    _, col, _ = st.columns([1, 2, 1])
    with col:
        _auth_header("Reset password", "We'll send you a reset link")

        if st.session_state.auth_error:
            st.error(st.session_state.auth_error)
        if st.session_state.auth_ok:
            st.success(st.session_state.auth_ok)

        email = st.text_input("Email address", key="rp_email",
                              placeholder="you@example.com")

        if st.button("Send reset link →", use_container_width=True,
                     type="primary", key="rp_btn"):
            st.session_state.auth_error = ""
            st.session_state.auth_ok    = ""
            if not email:
                st.session_state.auth_error = "Please enter your email address."
            else:
                try:
                    supabase.auth.reset_password_email(email)
                    st.session_state.auth_ok = (
                        "Reset link sent! Check your inbox and follow the link."
                    )
                except Exception as e:
                    st.session_state.auth_error = f"Failed to send reset email: {e}"
            st.rerun()

        st.markdown("<div style='height:.5rem'></div>", unsafe_allow_html=True)
        if st.button("← Back to sign in", use_container_width=True, key="rp_back"):
            st.session_state.auth_screen = "signin"
            st.session_state.auth_error  = ""
            st.session_state.auth_ok     = ""
            st.rerun()


# ─── Auth gate — show auth screens if not signed in ──────
if supabase is None:
    st.error("⚠️ Supabase is not configured. Add SUPABASE_URL and SUPABASE_KEY to your secrets.")
    st.stop()

if not st.session_state.user:
    if st.session_state.auth_screen == "signup":
        render_signup()
    elif st.session_state.auth_screen == "reset":
        render_reset()
    else:
        render_signin()
    st.stop()   # ← nothing below renders until signed in

# ─── Signed-in: load profile once per session ────────────
if not st.session_state.profile:
    st.session_state.profile   = db_get_profile(st.session_state.user.id)
    db_touch_last_seen(st.session_state.user.id)
    st.session_state.on_waitlist = db_is_on_waitlist(st.session_state.user.id)

# ─── Branding Header (only renders when signed in) ────────
initial = BRAND_NAME[0].upper() if BRAND_NAME else "?"
if BRAND_LOGO and os.path.exists(BRAND_LOGO):
    with open(BRAND_LOGO, "rb") as f:
        enc = base64.b64encode(f.read()).decode()
    _ext  = BRAND_LOGO.rsplit(".", 1)[-1].lower()
    _mime = "image/png" if _ext == "png" else "image/jpeg"
    logo_html = f'<img class="brand-logo" src="data:{_mime};base64,{enc}" />'
else:
    logo_html = f'<div class="brand-initials">{initial}</div>'

_display_name = ""
if st.session_state.profile:
    _display_name = st.session_state.profile.get("full_name", "") or st.session_state.user.email
else:
    _display_name = st.session_state.user.email

# Brand bar + user pill — sign out is a separate small button below
# Brand bar + sign out
_bar_col, _prem_col, _out_col = st.columns([5, 1, 1])
with _bar_col:
    st.markdown(f"""
    <div class="brand-bar" style="margin-bottom:0">
      <div class="brand-left">
        {logo_html}
        <div>
          <div class="brand-name">{BRAND_NAME}</div>
          <div class="brand-sub">{APP_SUBTITLE}</div>
        </div>
      </div>
      <div class="user-pill">
        <div class="user-pill-dot"></div>
        <span>{_display_name}</span>
      </div>
    </div>
    """, unsafe_allow_html=True)
with _prem_col:
    st.markdown("<div style='padding-top:1rem'></div>", unsafe_allow_html=True)
    _badge_label = "✓ Waitlisted" if st.session_state.on_waitlist else "⚡ Upgrade"
    if st.button(_badge_label, key="upgrade_btn", use_container_width=True):
        st.session_state.show_upgrade   = True
        st.session_state.upgrade_trigger = "manual"
with _out_col:
    st.markdown("<div style='padding-top:1rem'></div>", unsafe_allow_html=True)
    if st.button("↪ Sign out", key="signout_btn", use_container_width=True):
        try:
            supabase.auth.sign_out()
        except Exception:
            pass
        for _k, _v in {
            "user": None, "profile": None, "db": None,
            "summary": None, "active_doc_id": None,
            "processed_chunks": [], "chat_history": [],
            "summary_images": [], "summary_tables": [],
            "quiz_questions": [], "logs": [],
            "all_page_images": {}, "quiz_answers": {},
            "metrics": {"elements": 0, "chunks": 0, "docs": 0},
            "pipeline_ran": False, "quiz_submitted": False,
            "doc_name": "", "anon_uid": "",
            "auth_screen": "signin", "auth_error": "", "auth_ok": "",
        }.items():
            st.session_state[_k] = _v
        st.rerun()

# ─── MathJax — only injected after auth passes ───────────
st.markdown("""
<script>
window.MathJax = {
  tex: { inlineMath: [['$','$']], displayMath: [['$$','$$']] },
  svg: { fontCache: 'global' }
};
</script>
<script src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-svg.js" async></script>
""", unsafe_allow_html=True)

# ─── Hero ────────────────────────────────────────────────
st.markdown("""
<div class="hero-title">Ask your <span>documents.</span></div>
<div class="hero-desc">
Upload any document — PDF, Word, PowerPoint, Excel, images, and more.
The pipeline extracts text, tables &amp; images, builds a smart index,
then lets you chat with the content or test your understanding.
</div>
""", unsafe_allow_html=True)

# ─── Tabs ────────────────────────────────────────────────
tab_ingest, tab_summary, tab_query, tab_quiz, tab_logs, tab_premium = st.tabs(
    ["  📄  Upload & Index  ", "  📋  Summary  ", "  💬  Chat  ",
     "  🧠  Quiz  ", "  🗒  Logs  ", "  ⚡  Premium  "]
)

# ════════════════════════════════════════════════════════
# TAB 1 — INGEST
# ════════════════════════════════════════════════════════

# Supported types → (display label, icon, temp suffix)
SUPPORTED_TYPES = {
    "pdf":  ("PDF",        "📄"),
    "docx": ("Word Doc",   "📝"),
    "pptx": ("PowerPoint", "📊"),
    "xlsx": ("Excel",      "📗"),
    "txt":  ("Text File",  "📃"),
    "md":   ("Markdown",   "📃"),
    "html": ("HTML",       "🌐"),
    "csv":  ("CSV",        "📋"),
    "png":  ("Image",      "🖼️"),
    "jpg":  ("Image",      "🖼️"),
    "jpeg": ("Image",      "🖼️"),
}

with tab_ingest:

    # ── My Documents ─────────────────────────────────────
    user_docs = db_get_documents(st.session_state.user.id) if st.session_state.user else []

    if user_docs:
        st.markdown('<div class="sec-div"><hr/><span class="sec-lbl">My Documents</span><hr/></div>',
                    unsafe_allow_html=True)
        for doc in user_docs:
            is_active = st.session_state.active_doc_id == doc["id"]
            active_cls = "active" if is_active else ""
            created = doc.get("created_at", "")[:10]
            col_doc, col_del = st.columns([5, 1])
            with col_doc:
                st.markdown(f"""
                <div class="doc-card {active_cls}">
                    <div class="doc-card-name">📄 {doc['name']}</div>
                    <div class="doc-card-meta">{doc.get('file_type','').upper()} · 
                    {doc.get('chunk_count',0)} chunks · {created}</div>
                </div>
                """, unsafe_allow_html=True)
                if st.button("Load this document", key=f"load_{doc['id']}",
                             use_container_width=True):
                    # restore this document's chroma db into session
                    persist = doc.get("persist_dir", "")
                    if persist and os.path.exists(persist):
                        try:
                            from langchain_huggingface import HuggingFaceEmbeddings
                            from langchain_chroma import Chroma
                            embeddings = HuggingFaceEmbeddings(
                                model_name=DEFAULT_EMBED_MODEL,
                                model_kwargs={"device": "cpu"},
                            )
                            st.session_state.db = Chroma(
                                persist_directory=persist,
                                embedding_function=embeddings,
                                collection_metadata={"hnsw:space": "cosine"},
                            )
                            st.session_state.pipeline_ran  = True
                            st.session_state.doc_name      = doc["name"]
                            st.session_state.active_doc_id = doc["id"]
                            st.session_state.chat_history  = []
                            st.session_state.summary       = None
                            st.success(f"✅ Loaded: {doc['name']}")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Could not load document: {e}")
                    else:
                        st.warning("Document index not found on disk — please re-index.")
            with col_del:
                if st.button("🗑️", key=f"del_{doc['id']}",
                             help="Delete this document"):
                    db_delete_document(doc["id"])
                    if st.session_state.active_doc_id == doc["id"]:
                        st.session_state.pipeline_ran  = False
                        st.session_state.db            = None
                        st.session_state.active_doc_id = None
                    st.rerun()

        st.markdown('<div class="sec-div"><hr/><span class="sec-lbl">Index New Document</span><hr/></div>',
                    unsafe_allow_html=True)

    uploaded_file = st.file_uploader(
        "Upload document",
        type=list(SUPPORTED_TYPES.keys()),
        label_visibility="collapsed"
    )

    if uploaded_file:
        ext     = uploaded_file.name.rsplit(".", 1)[-1].lower()
        f_label, f_icon = SUPPORTED_TYPES.get(ext, ("Document", "📄"))
        size_kb = uploaded_file.size / 1024
        st.markdown(f"""
        <div class="file-pill">
            <div style="font-size:1.4rem">{f_icon}</div>
            <div>
                <div class="fp-name">{uploaded_file.name}</div>
                <div class="fp-size">{size_kb:.1f} KB · {f_label}</div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    run_btn = st.button(
        "Run Pipeline →",
        disabled=(uploaded_file is None or (_global_busy() and _busy_user() != _uid))
    )

    # ── busy notice for OTHER users ──────────────────────
    if _global_busy() and _busy_user() != _uid:
        elapsed  = _busy_elapsed_secs()
        remain   = max(0, PIPELINE_TIMEOUT_SECS - elapsed)
        mins_el  = elapsed // 60
        secs_el  = elapsed % 60
        mins_rem = remain // 60

        st.markdown(f"""
        <div class="busy-card">
            <div class="busy-spinner"></div>
            <div class="busy-title">Pipeline in use</div>
            <div class="busy-sub">
                Another user is currently indexing a document.<br>
                <strong>Running for {mins_el}m {secs_el:02d}s</strong>
                · auto-unlocks in ~{mins_rem} min if stuck.<br><br>
                You can still use the <strong>Chat</strong> and <strong>Quiz</strong> tabs
                if you've already indexed a document, or wait here and the button
                will unlock automatically when they're done.
            </div>
        </div>
        """, unsafe_allow_html=True)

        col_refresh, col_force = st.columns([3, 1])
        with col_force:
            if st.button("🔓 Force unlock", help="Use this if the pipeline has been stuck for more than 10 minutes"):
                _set_global_busy(False)
                st.success("Unlocked. You can now run the pipeline.")
                st.rerun()

        # auto-refresh every 8s so button unlocks without manual reload
        time.sleep(8)
        st.rerun()

    # supported formats strip
    st.markdown("""
    <div style="margin:.6rem 0 1.4rem 0; display:flex; flex-wrap:wrap; gap:6px; align-items:center;">
        <span style="font-size:.65rem; color:#525252; text-transform:uppercase; letter-spacing:.08em; margin-right:4px;">Supports</span>
        <span class="tag t-text">PDF</span>
        <span class="tag t-text">DOCX</span>
        <span class="tag t-text">PPTX</span>
        <span class="tag t-table">XLSX</span>
        <span class="tag t-table">CSV</span>
        <span class="tag t-text">TXT</span>
        <span class="tag t-text">MD</span>
        <span class="tag t-text">HTML</span>
        <span class="tag t-image">PNG</span>
        <span class="tag t-image">JPG</span>
    </div>
    """, unsafe_allow_html=True)

    # metrics + chunk preview after pipeline runs
    if st.session_state.pipeline_ran:
        m = st.session_state.metrics
        st.markdown(f"""
        <div class="metrics-row">
            <div class="metric-tile"><div class="metric-num">{m['elements']}</div><div class="metric-lbl">Elements</div></div>
            <div class="metric-tile"><div class="metric-num">{m['chunks']}</div><div class="metric-lbl">Chunks</div></div>
            <div class="metric-tile"><div class="metric-num">{m['docs']}</div><div class="metric-lbl">Indexed</div></div>
        </div>
        """, unsafe_allow_html=True)

        st.markdown('<div class="sec-div"><hr/><span class="sec-lbl">Chunk preview</span><hr/></div>', unsafe_allow_html=True)
        for i, doc in enumerate(st.session_state.processed_chunks[:5]):
            preview = doc.page_content[:250] + ("…" if len(doc.page_content) > 250 else "")
            orig    = json.loads(doc.metadata.get("original_content", "{}"))
            tags    = '<span class="tag t-text">text</span>'
            if orig.get("tables_html"):   tags += '<span class="tag t-table">table</span>'
            if orig.get("images_base64"): tags += '<span class="tag t-image">image</span>'
            with st.expander(f"Chunk {i + 1}"):
                st.markdown(f'<div style="margin-bottom:6px">{tags}</div>', unsafe_allow_html=True)
                st.markdown(f'<div class="chunk-card">{preview}</div>', unsafe_allow_html=True)

    # ── pipeline execution ──
    if run_btn and uploaded_file and not (_global_busy() and _busy_user() != _uid):
        st.session_state.logs = []
        _set_global_busy(True, _uid)  # 🔒 lock

        with st.status("Running pipeline…", expanded=True) as status:
            try:
                st.write("📂 Saving file…")
                ext      = uploaded_file.name.rsplit(".", 1)[-1].lower()
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                    tmp.write(uploaded_file.read())
                    tmp_path = tmp.name
                log(f"Saved → {tmp_path}")

                st.write("📦 Loading modules…")
                from unstructured.chunking.title import chunk_by_title
                from langchain_core.documents import Document
                from langchain_huggingface import HuggingFaceEmbeddings
                from langchain_chroma import Chroma
                log("Modules loaded", "success")

                # 1 ─ partition (route by file type) ──────────────────────
                f_label = SUPPORTED_TYPES.get(ext, ("Document", "📄"))[0]
                st.write(f"📄 Partitioning {f_label}…")

                elements = []

                if ext == "pdf":
                    # ── Tier 1: unstructured fast (best structured output) ──
                    try:
                        from unstructured.partition.pdf import partition_pdf
                        elements = partition_pdf(
                            filename=tmp_path,
                            strategy="fast",
                            infer_table_structure=True,
                        )
                        log(f"Tier 1 (unstructured fast): {len(elements)} elements")
                    except Exception as e:
                        log(f"Tier 1 failed: {e}", "error")
                        elements = []

                    # ── Tier 2: PyMuPDF text extraction ───────────────────
                    if not elements:
                        st.write("⚠️ Standard extraction got 0 elements — trying PyMuPDF text…")
                        try:
                            import fitz
                            from unstructured.documents.elements import Text, Title
                            doc_fitz = fitz.open(tmp_path)
                            for page in doc_fitz:
                                page_text = page.get_text("text").strip()
                                if page_text:
                                    # split into paragraphs on double newlines
                                    blocks = [b.strip() for b in page_text.split("\n\n") if b.strip()]
                                    for j, block in enumerate(blocks):
                                        # first block of each page treated as title candidate
                                        if j == 0 and len(block) < 120:
                                            el = Title(text=block)
                                        else:
                                            el = Text(text=block)
                                        el.metadata.page_number = page.number + 1
                                        elements.append(el)
                            doc_fitz.close()
                            log(f"Tier 2 (PyMuPDF text): {len(elements)} elements")
                            if elements:
                                st.write(f"✅ PyMuPDF extracted {len(elements)} text blocks")
                        except Exception as e:
                            log(f"Tier 2 failed: {e}", "error")
                            elements = []

                    # ── Tier 3: Tesseract OCR on page renders ─────────────
                    if not elements:
                        st.write("⚠️ No text layer found — running OCR on page images…")
                        try:
                            import fitz
                            import pytesseract
                            from PIL import Image as PILImage
                            import io
                            from unstructured.documents.elements import Text
                            doc_fitz  = fitz.open(tmp_path)
                            matrix    = fitz.Matrix(2.0, 2.0)   # 2x scale for better OCR
                            for page in doc_fitz:
                                pix      = page.get_pixmap(matrix=matrix, alpha=False)
                                img      = PILImage.open(io.BytesIO(pix.tobytes("png")))
                                ocr_text = pytesseract.image_to_string(img).strip()
                                if ocr_text:
                                    blocks = [b.strip() for b in ocr_text.split("\n\n") if b.strip()]
                                    for block in blocks:
                                        el = Text(text=block)
                                        el.metadata.page_number = page.number + 1
                                        elements.append(el)
                            doc_fitz.close()
                            log(f"Tier 3 (OCR): {len(elements)} elements")
                            if elements:
                                st.write(f"✅ OCR extracted {len(elements)} text blocks")
                        except Exception as e:
                            log(f"Tier 3 (OCR) failed: {e}", "error")

                elif ext == "docx":
                    from unstructured.partition.docx import partition_docx
                    elements = partition_docx(filename=tmp_path)
                elif ext == "pptx":
                    from unstructured.partition.pptx import partition_pptx
                    elements = partition_pptx(filename=tmp_path)
                elif ext == "xlsx":
                    from unstructured.partition.xlsx import partition_xlsx
                    elements = partition_xlsx(filename=tmp_path)
                elif ext in ("png", "jpg", "jpeg"):
                    # images: OCR directly
                    try:
                        import pytesseract
                        from PIL import Image as PILImage
                        from unstructured.documents.elements import Text
                        img      = PILImage.open(tmp_path)
                        ocr_text = pytesseract.image_to_string(img).strip()
                        if ocr_text:
                            for block in [b.strip() for b in ocr_text.split("\n\n") if b.strip()]:
                                elements.append(Text(text=block))
                    except Exception:
                        from unstructured.partition.image import partition_image
                        elements = partition_image(filename=tmp_path, strategy="fast")
                elif ext == "html":
                    from unstructured.partition.html import partition_html
                    elements = partition_html(filename=tmp_path)
                elif ext == "csv":
                    from unstructured.partition.csv import partition_csv
                    elements = partition_csv(filename=tmp_path)
                elif ext in ("txt", "md"):
                    from unstructured.partition.text import partition_text
                    elements = partition_text(filename=tmp_path)
                else:
                    from unstructured.partition.auto import partition
                    elements = partition(filename=tmp_path)

                # final check — if still empty after all fallbacks, abort cleanly
                if not elements:
                    raise ValueError(
                        "Could not extract any text from this document after trying "
                        "three methods (unstructured, PyMuPDF text layer, OCR). "
                        "The file may be corrupted, password-protected, or contain "
                        "only non-readable content."
                    )

                st.session_state.metrics["elements"] = len(elements)
                log(f"{len(elements)} elements extracted total", "success")
                st.write(f"✅ {len(elements)} elements extracted")

                # 1b ─ extract page images ────────────────────────────────
                st.write("🖼️ Extracting images…")
                page_images  = {}
                loose_images = []

                if ext == "pdf":
                    page_images = extract_images_from_pdf(tmp_path, dpi=150)
                    log(f"{len(page_images)} page renders captured", "success")
                elif ext == "docx":
                    loose_images = extract_images_from_docx(tmp_path)
                    log(f"{len(loose_images)} images from DOCX", "success")
                elif ext == "pptx":
                    loose_images = extract_images_from_pptx(tmp_path)
                    log(f"{len(loose_images)} images from PPTX", "success")

                if page_images or loose_images:
                    st.write(f"✅ {len(page_images) or len(loose_images)} image(s) captured")

                # 2 ─ chunk ───────────────────────────────────────────────
                st.write("🔨 Chunking…")
                try:
                    chunks = chunk_by_title(
                        elements,
                        max_characters=DEFAULT_MAX_CHARS,
                        new_after_n_chars=DEFAULT_NEW_AFTER,
                        combine_text_under_n_chars=DEFAULT_COMBINE,
                    )
                except Exception:
                    # chunk_by_title occasionally fails on synthetic elements
                    # fall back to simple fixed-size chunking
                    from langchain_core.documents import Document as LC_Doc
                    from unstructured.documents.elements import Text as UText
                    chunk_size = DEFAULT_MAX_CHARS
                    all_text   = "\n\n".join(el.text for el in elements if hasattr(el, "text") and el.text)
                    raw_chunks = [all_text[i:i+chunk_size] for i in range(0, len(all_text), chunk_size)]

                    # wrap as minimal objects with .text attribute
                    class _SimpleChunk:
                        def __init__(self, text, page_num=None):
                            self.text = text
                            class _Meta:
                                pass
                            self.metadata         = _Meta()
                            self.metadata.page_number    = page_num
                            self.metadata.orig_elements  = []
                    chunks = [_SimpleChunk(t) for t in raw_chunks if t.strip()]
                    log("Used simple fixed-size chunking fallback", "error")

                if page_images:
                    chunks = attach_page_images_to_chunks(chunks, page_images)

                st.session_state.metrics["chunks"] = len(chunks)
                log(f"{len(chunks)} chunks created", "success")
                st.write(f"✅ {len(chunks)} chunks")

                # 3 ─ AI summarise
                st.write("🧠 Generating AI summaries…")
                prog = st.progress(0)

                def separate(chunk, chunk_idx=0):
                    d = {"text": chunk.text, "tables": [], "images": []}

                    # ── resolve page numbers from chunk metadata ──────────
                    page_nums = set()
                    if hasattr(chunk, "_page_nums"):
                        page_nums = chunk._page_nums
                    elif hasattr(chunk, "metadata"):
                        if hasattr(chunk.metadata, "page_number") and chunk.metadata.page_number:
                            page_nums.add(int(chunk.metadata.page_number))
                        if hasattr(chunk.metadata, "orig_elements"):
                            for el in chunk.metadata.orig_elements:
                                if hasattr(el, "metadata") and hasattr(el.metadata, "page_number"):
                                    if el.metadata.page_number:
                                        page_nums.add(int(el.metadata.page_number))

                    # ── PDF: pull rendered page images by page number ──────
                    if page_images:
                        if page_nums:
                            for p in sorted(page_nums):
                                if p in page_images:
                                    d["images"].append(page_images[p])
                        else:
                            # no page number metadata — distribute proportionally
                            total_pages = len(page_images)
                            total_chunks = max(len(chunks), 1)
                            pages_per_chunk = max(1, total_pages // total_chunks)
                            start_page = chunk_idx * pages_per_chunk + 1
                            end_page   = start_page + pages_per_chunk
                            sorted_pages = sorted(page_images.keys())
                            for p in sorted_pages[start_page - 1 : end_page - 1]:
                                d["images"].append(page_images[p])

                    # ── DOCX/PPTX: distribute loose images evenly ─────────
                    if loose_images:
                        total_chunks = max(len(chunks), 1)
                        per_chunk    = max(1, len(loose_images) // total_chunks)
                        start        = chunk_idx * per_chunk
                        d["images"].extend(loose_images[start: start + per_chunk])

                    # ── tables from unstructured elements ─────────────────
                    if hasattr(chunk, "metadata") and hasattr(chunk.metadata, "orig_elements"):
                        for el in chunk.metadata.orig_elements:
                            if type(el).__name__ == "Table":
                                d["tables"].append(getattr(el.metadata, "text_as_html", el.text))

                    return d

                def ai_summary(text, tables, images):
                    try:
                        from langchain_core.messages import HumanMessage
                        p = f"Create a detailed, searchable description for retrieval.\n\nTEXT:\n{text}\n\n"
                        for i, t in enumerate(tables):
                            p += f"TABLE {i+1}:\n{t}\n\n"
                        p += "Cover key facts, numbers, topics, questions this answers, search terms, and describe any visible diagrams, figures, or formulas.\n\nDESCRIPTION:"
                        content = [{"type": "text", "text": p}]
                        # cap at 2 images per summary — page renders are large
                        for img in images[:2]:
                            # page renders are PNG; embedded images may be JPEG
                            mime = "image/png" if img.startswith("iVBOR") else "image/jpeg"
                            content.append({
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime};base64,{img}"}
                            })
                        response, _ = invoke_with_fallback([HumanMessage(content=content)])
                        return response.content
                    except Exception as e:
                        log(f"AI summary error: {e}", "error")
                        return text

                docs = []
                for i, chunk in enumerate(chunks):
                    cd = separate(chunk, chunk_idx=i)
                    enhanced = (
                        ai_summary(cd["text"], cd["tables"], cd["images"])
                        if (cd["tables"] or cd["images"]) else cd["text"]
                    )
                    docs.append(Document(
                        page_content=enhanced,
                        metadata={"original_content": json.dumps({
                            "raw_text":     cd["text"],
                            "tables_html":  cd["tables"],
                            "images_base64": cd["images"],
                        })}
                    ))
                    prog.progress((i + 1) / len(chunks))

                log(f"{len(docs)} docs processed", "success")
                st.write(f"✅ {len(docs)} docs processed")

                # 4 ─ vector store
                st.write("🔮 Building vector store…")
                embeddings = HuggingFaceEmbeddings(
                    model_name=DEFAULT_EMBED_MODEL,
                    model_kwargs={"device": "cpu"},
                )
                db = Chroma.from_documents(
                    documents=docs,
                    embedding=embeddings,
                    persist_directory=USER_PERSIST_DIR,
                    collection_metadata={"hnsw:space": "cosine"},
                )

                st.session_state.db               = db
                st.session_state.processed_chunks = docs
                st.session_state.metrics["docs"]  = len(docs)
                st.session_state.pipeline_ran      = True
                st.session_state.doc_name          = uploaded_file.name
                st.session_state.all_page_images   = page_images
                st.session_state.summary           = None
                log(f"Vector store ready → {USER_PERSIST_DIR}", "success")
                st.write(f"✅ {len(docs)} docs indexed")

                # ── save document record to Supabase ──────
                doc_id = db_save_document(
                    user_id     = st.session_state.user.id,
                    name        = uploaded_file.name,
                    file_type   = ext,
                    chunk_count = len(docs),
                    page_count  = len(page_images),
                    persist_dir = USER_PERSIST_DIR,
                )
                st.session_state.active_doc_id = doc_id
                log(f"Document saved to DB: {doc_id}", "success")

                # ── trigger upgrade nudge after 3rd document ──
                if not st.session_state.on_waitlist:
                    _doc_count = len(db_get_documents(st.session_state.user.id))
                    if _doc_count >= 3:
                        st.session_state.upgrade_trigger = "doc_limit"

                # ── open a fresh chat session ──────────────
                if doc_id:
                    sess_id = db_save_chat_session(st.session_state.user.id, doc_id)
                    st.session_state.chat_session_id = sess_id

                # 5 ─ generate summary
                st.write("📝 Generating document summary…")
                try:
                    st.session_state.summary = generate_summary(uploaded_file.name)
                    log("Summary generated", "success")
                    st.write("✅ Summary ready")
                except Exception as e:
                    log(f"Summary generation failed: {e}", "error")
                    st.write(f"⚠️ Summary failed (you can regenerate from the Summary tab): {e}")

                os.unlink(tmp_path)
                _set_global_busy(False)              # 🔓 unlock on success
                status.update(label="Pipeline complete ✅", state="complete")

            except Exception as e:
                log(f"Pipeline failed: {e}", "error")
                _set_global_busy(False)              # 🔓 unlock on error too
                status.update(label=f"Failed: {e}", state="error")
                st.error(str(e))

        st.rerun()

# ════════════════════════════════════════════════════════
# TAB 2 — SUMMARY
# ════════════════════════════════════════════════════════
with tab_summary:
    if not st.session_state.pipeline_ran:
        st.markdown("""
        <div style="text-align:center; padding:3rem 1rem;
                    border:1px dashed #2a2a2a; border-radius:14px; margin-top:1rem;">
            <div style="font-size:2.4rem; margin-bottom:.75rem">📋</div>
            <div style="font-family:'Syne',sans-serif; font-weight:700; font-size:1.05rem;
                        color:#f5f5f5; margin-bottom:.4rem;">No document indexed yet</div>
            <div style="color:#525252; font-size:.82rem">
                Upload and index a document first — the summary generates automatically.
            </div>
        </div>
        """, unsafe_allow_html=True)

    else:
        # ── regenerate button ────────────────────────────
        col_title, col_btn = st.columns([4, 1])
        with col_title:
            st.markdown(
                f'<div style="font-family:Syne,sans-serif;font-weight:800;'
                f'font-size:1.1rem;color:#f5f5f5;padding-top:.3rem;">'
                f'📋 {st.session_state.doc_name}</div>',
                unsafe_allow_html=True
            )
        with col_btn:
            if st.button("🔄 Regenerate", help="Re-generate the summary from scratch"):
                st.session_state.summary = None

        # ── generate if missing ──────────────────────────
        if st.session_state.summary is None:
            with st.spinner("Reading your document and generating summary… this may take a minute for long documents."):
                try:
                    st.session_state.summary = generate_summary(st.session_state.doc_name)
                    st.rerun()
                except Exception as e:
                    st.error(f"Summary failed: {e}. Click Regenerate to try again.")
                    st.stop()

        s = st.session_state.summary

        # ════ HERO — what is this document? ══════════════
        st.markdown(f"""
        <div class="sum-hero">
            <div class="sum-doc-type">📄 Document Summary</div>
            <div class="sum-doc-title">{s.get("topic", st.session_state.doc_name)}</div>
            <div class="sum-plain-english">{s.get("plain_english", "").replace(chr(10), "<br>")}</div>
        </div>
        """, unsafe_allow_html=True)

        # ════ KEY FIGURES ═════════════════════════════════
        key_imgs = st.session_state.summary_images
        if key_imgs:
            st.markdown('<div class="sum-section"><div class="sum-section-title">🖼️ Key Figures from Document</div>', unsafe_allow_html=True)
            cols = st.columns(min(len(key_imgs), 3))
            for i, b64 in enumerate(key_imgs):
                try:
                    cols[i % 3].image(base64.b64decode(b64), use_container_width=True)
                    cols[i % 3].markdown(f'<div class="sum-fig-caption">Figure {i+1}</div>', unsafe_allow_html=True)
                except Exception:
                    pass
            st.markdown('</div>', unsafe_allow_html=True)

        # ════ SECTION-BY-SECTION BREAKDOWN ═══════════════
        sections = s.get("sections", [])
        if sections:
            st.markdown('<div class="sum-section"><div class="sum-section-title">📖 Full Document Breakdown</div>', unsafe_allow_html=True)
            for sec in sections:
                st.markdown(f"### {sec.get('title', 'Section')}")
                st.markdown(sec.get("summary", ""))

                # inline formulas for this section
                sec_formulas = sec.get("formulas", [])
                if sec_formulas:
                    for f in sec_formulas:
                        st.markdown(f)

                # key point callout
                kp = sec.get("key_point", "")
                if kp:
                    st.markdown(f"""
                    <div style="background:rgba(249,115,22,.07);border-left:3px solid var(--orange);
                                border-radius:0 8px 8px 0;padding:.6rem 1rem;margin:.6rem 0;
                                font-size:.83rem;color:#fb923c;">
                        💡 <strong>Key point:</strong> {kp}
                    </div>
                    """, unsafe_allow_html=True)

                st.markdown("<hr style='border-color:#2a2a2a;margin:1rem 0'>", unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        # ════ KEY CONCEPTS GLOSSARY ═══════════════════════
        concepts = s.get("concepts", [])
        if concepts:
            st.markdown('<div class="sum-section"><div class="sum-section-title">🔑 Key Concepts & Terms</div>', unsafe_allow_html=True)
            for c in concepts:
                st.markdown(f"""
                <div class="sum-concept-card">
                    <div class="sum-concept-term">{c.get("term", "")}</div>
                    <div class="sum-concept-def">{c.get("definition", "")}</div>
                </div>
                """, unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        # ════ FORMULAS ════════════════════════════════════
        formulas = s.get("formulas", [])
        if formulas:
            st.markdown('<div class="sum-section"><div class="sum-section-title">🧮 Formulas & Equations</div>', unsafe_allow_html=True)
            for f in formulas:
                label = f.get("label", "")
                latex = f.get("latex", "")
                expl  = f.get("explanation", "")
                if label:
                    st.markdown(f"**{label}**")
                if latex:
                    st.markdown(latex)
                if expl:
                    st.markdown(f'<div style="font-size:.82rem;color:#a3a3a3;margin-bottom:.8rem;">{expl}</div>', unsafe_allow_html=True)
                st.markdown("<div style='height:.3rem'></div>", unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        # ════ TABLES ══════════════════════════════════════
        tables = s.get("tables", [])
        if tables:
            st.markdown('<div class="sum-section"><div class="sum-section-title">📊 Tables</div>', unsafe_allow_html=True)
            for t in tables:
                title = t.get("title", "")
                md    = t.get("markdown", "")
                if title:
                    st.markdown(f"**{title}**")
                if md:
                    st.markdown(md)
                st.markdown("<div style='height:.5rem'></div>", unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        # ════ TAKEAWAYS ═══════════════════════════════════
        takeaways = s.get("takeaways", [])
        if takeaways:
            st.markdown('<div class="sum-section"><div class="sum-section-title">💡 What You Must Remember</div>', unsafe_allow_html=True)
            for i, t in enumerate(takeaways):
                st.markdown(f"""
                <div class="sum-takeaway">
                    <div class="sum-takeaway-num">{i+1}</div>
                    <div>{t}</div>
                </div>
                """, unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        # ── bottom CTA ───────────────────────────────────
        st.markdown("""
        <div style="text-align:center;padding:1.5rem 0 .5rem 0;color:#525252;font-size:.8rem;">
            Still have questions? Switch to the <strong style="color:#a3a3a3">Chat</strong> tab
            to ask anything, or test yourself in the <strong style="color:#a3a3a3">Quiz</strong> tab.
        </div>
        """, unsafe_allow_html=True)


# ════════════════════════════════════════════════════════
# TAB 3 — CHAT
# ════════════════════════════════════════════════════════
with tab_query:
    if not st.session_state.pipeline_ran:
        st.markdown("""
        <div style="text-align:center; padding:3rem 1rem;
                    border:1px dashed #2a2a2a; border-radius:14px; margin-top:1rem;">
            <div style="font-size:2.4rem; margin-bottom:.75rem">💬</div>
            <div style="font-family:'Syne',sans-serif; font-weight:700; font-size:1.05rem;
                        color:#f5f5f5; margin-bottom:.4rem;">No document indexed yet</div>
            <div style="color:#525252; font-size:.82rem">
                Go to <strong style="color:#a3a3a3">Upload &amp; Index</strong> to get started
            </div>
        </div>
        """, unsafe_allow_html=True)
    else:
        # replay history
        for turn in st.session_state.chat_history:
            with st.chat_message("user"):
                st.write(turn["query"])
            with st.chat_message("assistant"):
                atype = turn.get("answer_type", "doc")
                render_answer(turn["answer"], turn.get("images", []), is_gk=(atype == "gk"))
                if turn.get("chunks"):
                    with st.expander(f"📎 {len(turn['chunks'])} source chunks"):
                        for i, c in enumerate(turn["chunks"]):
                            st.markdown(
                                f'<div class="chunk-card"><strong>Source {i+1}</strong>'
                                f'<br>{c.page_content[:300]}…</div>',
                                unsafe_allow_html=True
                            )

        query = st.chat_input("Ask anything about your document…")

        if query:
            with st.chat_message("user"):
                st.write(query)

            with st.chat_message("assistant"):
                with st.spinner("Thinking…"):
                    try:
                        from langchain_core.messages import HumanMessage

                        # ── retrieve with similarity scores ──────────────
                        results_with_scores = st.session_state.db.similarity_search_with_score(
                            query, k=DEFAULT_TOP_K
                        )
                        retrieved = [r[0] for r in results_with_scores]
                        scores    = [r[1] for r in results_with_scores]

                        # ── deduplicate content ──────────────────────────
                        chunk_images, chunk_tables, chunk_texts = collect_content(retrieved)

                        # ── decide: doc answer, GK, or both ─────────────
                        use_gk, gk_reason = should_use_general_knowledge(query, retrieved, scores)

                        # ─── PROMPT BUILDER ──────────────────────────────
                        FORMAT_RULES = """
Formatting rules — follow these strictly:
- Use $...$ for inline math (e.g. $x^2$)
- Use $$...$$ on its own line for block equations
- Use markdown tables for tabular data
- Use **bold** for key terms, `code` for variable names
- Use bullet or numbered lists where appropriate
- Never write raw LaTeX without $ signs
"""
                        notice_slot = st.empty()

                        # ── Path A: answer from document ─────────────────
                        if not use_gk or gk_reason == "intent":
                            doc_prompt  = f"Answer this question using ONLY the documents below.\n\nQUESTION: {query}\n\nDOCUMENTS:\n"
                            for i, txt in enumerate(chunk_texts):
                                doc_prompt += f"\n--- Text block {i+1} ---\n{txt}\n"
                            if chunk_tables:
                                doc_prompt += "\n--- Tables ---\n"
                                for j, tbl in enumerate(chunk_tables):
                                    doc_prompt += f"Table {j+1}:\n{tbl}\n\n"
                            if chunk_images:
                                doc_prompt += f"\n{len(chunk_images)} document image(s) attached — reference them where relevant.\n"
                            doc_prompt += FORMAT_RULES + "\nProvide a clear, complete answer from the document. If the document does not contain enough information, say so explicitly.\n\nANSWER:"

                            content = [{"type": "text", "text": doc_prompt}]
                            for b64 in chunk_images[:4]:
                                mime = "image/png" if b64.startswith("iVBOR") else "image/jpeg"
                                content.append({"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}})

                            doc_response, provider = invoke_with_fallback([HumanMessage(content=content)], status_slot=notice_slot)
                            doc_answer = doc_response.content

                            if provider == "groq":
                                notice_slot.markdown('<div style="font-size:.65rem;color:#525252;margin-bottom:6px;">⚡ Groq · Llama 3.3 70B (images not analysed this turn)</div>', unsafe_allow_html=True)
                            else:
                                notice_slot.empty()

                            render_answer(doc_answer, chunk_images, is_gk=False)
                            answer      = doc_answer
                            answer_type = "doc"

                        # ── Path B: pure general knowledge ───────────────
                        if use_gk and gk_reason in ("low_confidence", "both"):
                            gk_prompt = f"""You are a knowledgeable tutor. The user's document did not contain a good answer to this question, so answer from your general knowledge.

QUESTION: {query}

{FORMAT_RULES}

Give a thorough, clear explanation with examples and analogies where helpful. Be educational.

ANSWER:"""
                            gk_response, _ = invoke_with_fallback([HumanMessage(content=gk_prompt)])
                            gk_answer       = gk_response.content
                            render_answer(gk_answer, [], is_gk=True)
                            answer      = gk_answer
                            answer_type = "gk"

                        # ── Path C: doc answer + GK expansion (intent match) ──
                        if use_gk and gk_reason == "intent":
                            gk_expand_prompt = f"""The user asked: "{query}"

A document-based answer was already given. Now give a concise general-knowledge explanation of the core concept(s) involved, as a tutor would — using examples and analogies. Keep it brief (3-5 sentences).

{FORMAT_RULES}

GENERAL EXPLANATION:"""
                            gk_response, _ = invoke_with_fallback([HumanMessage(content=gk_expand_prompt)])
                            gk_answer       = gk_response.content
                            st.markdown('<div class="sec-div"><hr/><span class="sec-lbl">General knowledge expansion</span><hr/></div>', unsafe_allow_html=True)
                            render_answer(gk_answer, [], is_gk=True)
                            answer_type = "hybrid"

                        with st.expander(f"📎 {len(retrieved)} source chunks · best score: {min(scores):.3f}"):
                            for i, (c, s) in enumerate(zip(retrieved, scores)):
                                st.markdown(
                                    f'<div class="chunk-card"><strong>Source {i+1}</strong> '
                                    f'<span style="color:{"#22c55e" if s < 0.4 else "#f97316" if s < 0.65 else "#ef4444"};font-size:.7rem;">score {s:.3f}</span>'
                                    f'<br>{c.page_content[:300]}…</div>',
                                    unsafe_allow_html=True
                                )

                        st.session_state.chat_history.append({
                            "query":       query,
                            "answer":      answer,
                            "answer_type": answer_type,
                            "chunks":      retrieved,
                            "scores":      scores,
                            "images":      chunk_images,
                        })

                        # ── persist to Supabase ───────────
                        sess_id = st.session_state.get("chat_session_id")
                        if sess_id:
                            db_save_message(sess_id, "user",      query,  "user")
                            db_save_message(sess_id, "assistant", answer, answer_type)

                        # ── trigger nudge after 10th message ──
                        if not st.session_state.on_waitlist:
                            _msg_count = len(st.session_state.chat_history)
                            if _msg_count >= 10 and not st.session_state.upgrade_trigger:
                                st.session_state.upgrade_trigger = "chat_limit"

                    except Exception as e:
                        st.error(f"Error: {e}")

        if st.session_state.chat_history:
            if st.button("🗑 Clear chat"):
                st.session_state.chat_history = []
                st.rerun()

# ════════════════════════════════════════════════════════
# TAB 4 — QUIZ
# ════════════════════════════════════════════════════════
with tab_quiz:
    if not st.session_state.pipeline_ran:
        st.markdown("""
        <div style="text-align:center; padding:3rem 1rem;
                    border:1px dashed #2a2a2a; border-radius:14px; margin-top:1rem;">
            <div style="font-size:2.4rem; margin-bottom:.75rem">🧠</div>
            <div style="font-family:'Syne',sans-serif; font-weight:700; font-size:1.05rem;
                        color:#f5f5f5; margin-bottom:.4rem;">No document indexed yet</div>
            <div style="color:#525252; font-size:.82rem">
                Upload and index a document first, then come back here to test your understanding.
            </div>
        </div>
        """, unsafe_allow_html=True)
    else:
        # ── setup screen ─────────────────────────────────
        if not st.session_state.quiz_questions:
            st.markdown("""
            <div style="margin-bottom:1.5rem;">
                <div style="font-family:'Syne',sans-serif; font-weight:800; font-size:1.6rem;
                            letter-spacing:-0.03em; color:#f5f5f5; margin-bottom:.3rem;">
                    Test your understanding
                </div>
                <div style="color:#a3a3a3; font-size:.85rem; line-height:1.5;">
                    The quiz is generated directly from your uploaded document.
                    Pick a difficulty and number of questions to begin.
                </div>
            </div>
            """, unsafe_allow_html=True)

            col_a, col_b = st.columns(2)
            with col_a:
                difficulty = st.selectbox(
                    "Difficulty",
                    ["Easy", "Medium", "Hard"],
                    index=1,
                    label_visibility="visible"
                )
            with col_b:
                num_q = st.selectbox(
                    "Number of questions",
                    [3, 5, 10],
                    index=1,
                    label_visibility="visible"
                )

            diff_class = {"Easy": "easy", "Medium": "medium", "Hard": "hard"}[difficulty]
            st.markdown(f"""
            <div style="margin:.8rem 0 1.2rem 0;">
                <span class="diff-badge diff-{diff_class}">{difficulty}</span>
                <span style="color:#525252; font-size:.78rem;">
                    {num_q} questions · generated from your document
                </span>
            </div>
            """, unsafe_allow_html=True)

            if st.button("🧠 Generate Quiz"):
                with st.spinner("Generating questions from your document…"):
                    try:
                        qs = generate_quiz(num_q, difficulty)
                        st.session_state.quiz_questions = qs
                        st.session_state.quiz_answers   = {}
                        st.session_state.quiz_submitted  = False
                        st.rerun()
                    except json.JSONDecodeError:
                        st.error("The AI returned malformed JSON. Please try again.")
                    except Exception as e:
                        st.error(f"Quiz generation failed: {e}")

        # ── active quiz ───────────────────────────────────
        elif not st.session_state.quiz_submitted:
            questions = st.session_state.quiz_questions
            answered  = len(st.session_state.quiz_answers)
            total     = len(questions)

            # progress bar
            st.progress(answered / total)
            st.markdown(
                f'<div style="font-size:.72rem; color:#525252; margin-bottom:1.2rem;">'
                f'{answered} of {total} answered</div>',
                unsafe_allow_html=True
            )

            with st.form("quiz_form"):
                user_answers = {}
                for i, q in enumerate(questions):
                    st.markdown(f"""
                    <div class="quiz-q">
                        <div class="quiz-q-num">Question {i + 1} of {total}</div>
                        <div class="quiz-q-text">{q['question']}</div>
                    </div>
                    """, unsafe_allow_html=True)

                    choice = st.radio(
                        f"q_{i}",
                        options=q["options"],
                        index=st.session_state.quiz_answers.get(i, None),
                        label_visibility="collapsed",
                        key=f"radio_{i}"
                    )
                    user_answers[i] = q["options"].index(choice) if choice else None
                    st.markdown("<div style='margin-bottom:1rem'></div>", unsafe_allow_html=True)

                submitted = st.form_submit_button("✅ Submit Answers")
                if submitted:
                    if any(v is None for v in user_answers.values()):
                        st.warning("Please answer all questions before submitting.")
                    else:
                        st.session_state.quiz_answers  = user_answers
                        st.session_state.quiz_submitted = True
                        # ── save quiz result to Supabase ──
                        total_q   = len(st.session_state.quiz_questions)
                        correct_q = sum(
                            1 for i, q in enumerate(st.session_state.quiz_questions)
                            if user_answers.get(i) == q["answer"]
                        )
                        db_save_quiz(
                            user_id     = st.session_state.user.id,
                            document_id = st.session_state.get("active_doc_id"),
                            difficulty  = st.session_state.get("quiz_difficulty", "medium"),
                            score       = correct_q,
                            total       = total_q,
                        )
                        st.rerun()

        # ── results screen ────────────────────────────────
        else:
            questions = st.session_state.quiz_questions
            answers   = st.session_state.quiz_answers
            total     = len(questions)
            correct   = sum(1 for i, q in enumerate(questions) if answers.get(i) == q["answer"])
            pct       = round((correct / total) * 100)

            # score card
            if pct >= 80:
                grade_color = "#22c55e"
                grade_msg   = "Excellent work! 🎉"
            elif pct >= 50:
                grade_color = "#f97316"
                grade_msg   = "Good effort — review the sections below."
            else:
                grade_color = "#ef4444"
                grade_msg   = "Keep studying — you'll get there!"

            st.markdown(f"""
            <div class="quiz-score-card">
                <div class="quiz-score-num" style="color:{grade_color};">{correct}/{total}</div>
                <div class="quiz-score-lbl">correct answers · {pct}%</div>
                <div style="color:{grade_color}; font-size:.9rem; margin-top:.6rem; font-weight:500;">
                    {grade_msg}
                </div>
            </div>
            """, unsafe_allow_html=True)

            # per-question review
            st.markdown('<div class="sec-div"><hr/><span class="sec-lbl">Answer review</span><hr/></div>', unsafe_allow_html=True)

            for i, q in enumerate(questions):
                user_ans    = answers.get(i)
                correct_ans = q["answer"]
                is_correct  = user_ans == correct_ans

                icon = "✅" if is_correct else "❌"
                with st.expander(f"{icon}  Q{i+1}: {q['question'][:80]}…"):
                    for j, opt in enumerate(q["options"]):
                        if j == correct_ans and j == user_ans:
                            tag = "✅ Your answer · Correct"
                            css = "correct"
                        elif j == correct_ans:
                            tag = "✅ Correct answer"
                            css = "correct"
                        elif j == user_ans:
                            tag = "❌ Your answer"
                            css = "wrong"
                        else:
                            tag = ""
                            css = ""

                        border = f"border-color:#22c55e" if css == "correct" else (f"border-color:#ef4444" if css == "wrong" else "")
                        color  = f"color:#22c55e" if css == "correct" else (f"color:#ef4444" if css == "wrong" else "color:#a3a3a3")
                        st.markdown(f"""
                        <div class="quiz-opt" style="{border}; {color};">
                            <span style="font-weight:600; min-width:18px">{"ABCD"[j]}.</span>
                            <span>{opt}</span>
                            {f'<span style="margin-left:auto; font-size:.7rem; opacity:.8">{tag}</span>' if tag else ''}
                        </div>
                        """, unsafe_allow_html=True)

                    if q.get("explanation"):
                        st.markdown(f'<div class="quiz-explanation">💡 {q["explanation"]}</div>', unsafe_allow_html=True)

            st.markdown("<div style='height:1rem'></div>", unsafe_allow_html=True)
            col1, col2 = st.columns(2)
            with col1:
                if st.button("🔄 Retake this quiz"):
                    st.session_state.quiz_answers   = {}
                    st.session_state.quiz_submitted  = False
                    st.rerun()
            with col2:
                if st.button("🆕 New quiz"):
                    st.session_state.quiz_questions = []
                    st.session_state.quiz_answers   = {}
                    st.session_state.quiz_submitted  = False
                    st.rerun()

# ════════════════════════════════════════════════════════
# TAB 5 — LOGS
# ════════════════════════════════════════════════════════
with tab_logs:
    st.markdown("### Pipeline logs")
    if not st.session_state.logs:
        st.markdown('<div style="color:#525252; font-size:.82rem; padding:.5rem 0">No logs yet.</div>', unsafe_allow_html=True)
    else:
        lines = ""
        for e in st.session_state.logs:
            cls = "ok" if e["level"] == "success" else ("err" if e["level"] == "error" else "")
            px  = "✅" if e["level"] == "success" else ("❌" if e["level"] == "error" else "›")
            lines += f'<div class="ll {cls}"><span class="px">{px}</span>{e["msg"]}</div>'
        st.markdown(f'<div class="log-wrap">{lines}</div>', unsafe_allow_html=True)

        if st.button("🗑 Clear logs"):
            st.session_state.logs = []
            st.rerun()


# ════════════════════════════════════════════════════════
# TAB 6 — PREMIUM / WAITLIST
# ════════════════════════════════════════════════════════
with tab_premium:

    _wl_user  = st.session_state.user
    _wl_prof  = st.session_state.profile or {}
    _wl_name  = _wl_prof.get("full_name", "") or ""
    _wl_email = _wl_prof.get("email", "") or _wl_user.email
    _wl_count = db_waitlist_count()

    # ── already on waitlist ───────────────────────────────
    if st.session_state.on_waitlist:
        st.markdown("""
        <div class="waitlist-hero">
            <div class="waitlist-hero-icon">🎉</div>
            <div class="waitlist-hero-title">You're on the list!</div>
            <div class="waitlist-hero-sub">
                We'll email you the moment Premium launches.<br>
                You'll be among the very first to get access.
            </div>
        </div>
        """, unsafe_allow_html=True)

        st.markdown(f"""
        <div class="waitlist-joined">
            <div class="waitlist-joined-title">✓ Waitlist confirmed</div>
            <div class="waitlist-joined-sub">
                {_wl_count} people are waiting for Premium · We'll be in touch soon
            </div>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("---")
        st.markdown("#### While you wait, here's what's coming:")

    # ── not yet on waitlist ───────────────────────────────
    else:
        # nudge banner if a trigger fired
        if st.session_state.upgrade_trigger == "doc_limit":
            st.markdown("""
            <div class="upgrade-nudge">
                <div class="upgrade-nudge-text">
                    <div class="upgrade-nudge-title">You've indexed 3 documents 🔥</div>
                    <div class="upgrade-nudge-sub">
                        Serious users like you are exactly who Premium is built for.
                        Join the waitlist below to lock in early access.
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)
        elif st.session_state.upgrade_trigger == "chat_limit":
            st.markdown("""
            <div class="upgrade-nudge">
                <div class="upgrade-nudge-text">
                    <div class="upgrade-nudge-title">10 messages and counting 💬</div>
                    <div class="upgrade-nudge-sub">
                        You're clearly getting value from this. Premium gives you
                        unlimited chats, bigger documents, and a lot more.
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)

        st.markdown("""
        <div class="waitlist-hero">
            <div class="waitlist-hero-icon">⚡</div>
            <div class="waitlist-hero-title">Go Premium — Early Access</div>
            <div class="waitlist-hero-sub">
                The free version is powerful. Premium is everything else.<br>
                Join the waitlist and be first to know when it launches.
            </div>
        </div>
        """, unsafe_allow_html=True)

        if _wl_count > 0:
            st.markdown(
                f'<div style="text-align:center; font-size:.8rem; color:#f97316; '
                f'margin-bottom:1rem; font-weight:600;">'
                f'🔥 {_wl_count} people already on the waitlist</div>',
                unsafe_allow_html=True
            )

    # ── feature list (always shown) ───────────────────────
    st.markdown("""
    <div style="background:#1a1a1a; border:1px solid #2a2a2a;
                border-radius:12px; padding:1.2rem 1.4rem; margin-bottom:1.5rem;">
    """, unsafe_allow_html=True)

    features = [
        ("🚀", "Faster processing",
         "<strong>Priority queue</strong> — your documents never wait behind other users"),
        ("📄", "Bigger documents",
         "Process files up to <strong>500 pages</strong> (free tier: 50 pages)"),
        ("💬", "Unlimited chat",
         "<strong>No message limits</strong> — ask as many questions as you need"),
        ("🧠", "Smarter model",
         "Answers powered by <strong>GPT-4o</strong> with higher accuracy on complex content"),
        ("📚", "Document library",
         "Store and switch between <strong>unlimited documents</strong> in your library"),
        ("🧩", "Quiz packs",
         "Save quizzes, track your scores over time, and <strong>share with friends</strong>"),
        ("📤", "Export answers",
         "Download your chat history and summaries as <strong>PDF or Word</strong>"),
        ("🛟", "Priority support",
         "Direct access to the team — <strong>responses within 24 hours</strong>"),
    ]
    for icon, title, desc in features:
        st.markdown(f"""
        <div class="feature-row">
            <div class="feature-icon">{icon}</div>
            <div class="feature-text"><strong>{title}</strong> — {desc}</div>
        </div>
        """, unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)

    # ── waitlist form (only if not already on it) ─────────
    if not st.session_state.on_waitlist:
        st.markdown("### Join the waitlist")
        st.markdown(
            '<div style="font-size:.82rem; color:#a3a3a3; margin-bottom:1rem;">'
            'Takes 10 seconds. No payment needed — just your interest.</div>',
            unsafe_allow_html=True
        )

        _wl_name_input  = st.text_input("Your name",
                                         value=_wl_name, key="wl_name")
        _wl_email_input = st.text_input("Email address",
                                         value=_wl_email, key="wl_email")

        st.markdown(
            '<div style="font-size:.82rem; color:#a3a3a3; margin:.6rem 0 .3rem 0;">'
            'How much would you pay per month for Premium?</div>',
            unsafe_allow_html=True
        )
        _price = st.radio(
            "Price range",
            ["Not sure yet", "$5 – $10 / month", "$10 – $20 / month", "$20+ / month"],
            key="wl_price",
            label_visibility="collapsed",
            horizontal=True,
        )

        _wl_error = ""
        if st.button("⚡ Join the waitlist", type="primary",
                     use_container_width=True, key="wl_submit"):
            if not _wl_name_input.strip() or not _wl_email_input.strip():
                _wl_error = "Please fill in your name and email."
            else:
                _ok = db_join_waitlist(
                    user_id     = _wl_user.id,
                    full_name   = _wl_name_input.strip(),
                    email       = _wl_email_input.strip(),
                    trigger     = st.session_state.upgrade_trigger or "manual",
                    price_range = _price,
                )
                if _ok:
                    st.session_state.on_waitlist     = True
                    st.session_state.upgrade_trigger = ""
                    st.rerun()
                else:
                    _wl_error = "Something went wrong — please try again."

        if _wl_error:
            st.error(_wl_error)

        st.markdown(
            '<div style="font-size:.72rem; color:#525252; text-align:center; '
            'margin-top:.8rem;">No spam. No payment. Just an early heads-up.</div>',
            unsafe_allow_html=True
        )


# ── Upgrade nudge banner (shown inline in ingest/chat tabs after triggers) ──
# Rendered outside tabs so it appears as a floating reminder at page bottom
if st.session_state.upgrade_trigger and not st.session_state.on_waitlist:
    _trig_msg = {
        "doc_limit":  "You've indexed 3 documents — you're a power user. Want more?",
        "chat_limit": "10 messages sent — you're getting real value. Want unlimited?",
        "manual":     "Interested in more? Check out Premium.",
    }.get(st.session_state.upgrade_trigger, "")

    if _trig_msg:
        st.markdown(f"""
        <div class="upgrade-nudge" style="margin-top:1.5rem;">
            <div class="upgrade-nudge-text">
                <div class="upgrade-nudge-title">{_trig_msg}</div>
                <div class="upgrade-nudge-sub">
                    Open the ⚡ Premium tab to see what's coming and join the waitlist.
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)
